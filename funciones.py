# import pandas as pd
# import plotly.express as px
# import plotly.graph_objects as go
# import streamlit as st

# def graficar_barras_porcentaje(df, col_x, col_color, col_valores="Personas", titulo="Avance", lista_colores=None):
#     """
#     Crea una gráfica de barras apiladas en porcentajes optimizada para Streamlit.
#     Soporta una lista de colores personalizada de cualquier tamaño.
#     """
#     # 1. Validación de seguridad para que la app nunca falle
#     if df is None or df.empty or col_x not in df.columns or col_color not in df.columns:
#         fig_vacia = go.Figure()
#         fig_vacia.update_layout(
#             title=f"{titulo} (Sin datos disponibles)",
#             xaxis={"visible": False}, yaxis={"visible": False},
#             annotations=[{"text": "No hay datos para mostrar", "textangle": 0, "showarrow": False, "font": {"size": 20}}]
#         )
#         return fig_vacia

#     try:
#         # 2. Agrupar y calcular totales/porcentajes
#         df_agrupado = df.groupby([col_x, col_color])[col_valores].sum().reset_index()
#         df_total = df_agrupado.groupby(col_x)[col_valores].sum().reset_index().rename(columns={col_valores: 'Total_X'})
        
#         df_procesado = df_agrupado.merge(df_total, on=col_x)
#         df_procesado["Porcentaje"] = (df_procesado[col_valores] / df_procesado["Total_X"] * 100).round(1)

#         # 3. Definir orden descendente basado en la primera categoría del color
#         categorias_color = df_procesado[col_color].unique()
#         primer_color = categorias_color[0] if len(categorias_color) > 0 else ""
        
#         df_orden = df_procesado[df_procesado[col_color] == primer_color].sort_values("Porcentaje", ascending=False)
#         orden_x = df_orden[col_x].tolist()
        
#         todos_x = df_procesado[col_x].unique()
#         orden_final = orden_x + [cat for cat in todos_x if cat not in orden_x]

#         # 4. Configuración dinámica del grosor de barras
#         bar_gap = 0.8 if df_procesado[col_x].nunique() <= 3 else 0.15

#         # 5. ASIGNACIÓN INTELIGENTE DE COLORES
#         color_map = {}
#         if lista_colores and isinstance(lista_colores, list) and len(lista_colores) > 0:
#             # Toma solo los colores necesarios. Si faltan, usa operador residuo (%) para ciclar la lista
#             color_map = {cat: lista_colores[i % len(lista_colores)] for i, cat in enumerate(categorias_color)}
#         else:
#             # Paleta institucional fija por defecto si los nombres coinciden, o paleta segura de Plotly
#             color_map_base = {"actualizadas": "#235B4E", "pendientes": "#D4C19C", "Si": "#235B4E", "No": "#D4C19C"}
#             color_map = {cat: color_map_base.get(cat, px.colors.qualitative.Safe[i % 10]) for i, cat in enumerate(categorias_color)}

#         # 6. Construcción de la gráfica
#         fig = px.bar(
#             df_procesado, x=col_x, y="Porcentaje", color=col_color,
#             text=df_procesado["Porcentaje"].apply(lambda x: f"{x:.1f}%"),
#             custom_data=[col_x, col_valores, "Porcentaje", col_color],
#             color_discrete_map=color_map,
#             category_orders={col_x: orden_final},
#             barmode="stack",
#         )

#         # 7. Estilizado de textos internos y Tooltips
#         fig.update_traces(
#             textposition='inside', insidetextanchor='middle',
#             textfont=dict(size=14, color="white"),
#             hovertemplate="<b>%{customdata[0]}:</b><br> %{customdata[1]:,.0f} (%{customdata[2]:.1f}%)<br>%{customdata[3]} <extra></extra>"
#         )

#         # 8. Diseño del Layout
#         fig.update_layout(
#             title=dict(text=titulo, font=dict(size=20, color="#6A1B29")), 
#             xaxis=dict(title=dict(text=str(col_x).capitalize()), tickangle=-45),
#             yaxis=dict(title=dict(text="Porcentaje (%)"), range=[0, 115], ticksuffix="%", dtick=20),
#             bargap=bar_gap, height=600, template="plotly_white",
#             legend=dict(title=dict(text=str(col_color).capitalize()), orientation="h", yanchor="bottom", y=1.02, xanchor="center", x=0.5),
#             margin=dict(l=50, t=80, r=30, b=130), uniformtext_minsize=9, uniformtext_mode='hide', paper_bgcolor="rgba(0,0,0,0)"
#         )
#         return fig

#     except Exception:
#         fig_error = go.Figure()
#         fig_error.update_layout(title="Error al procesar la estructura de datos.", xaxis={"visible": False}, yaxis={"visible": False})
#         return fig_error


# import streamlit as st
# import pandas as pd

# # --- EJEMPLO 1: Con tus datos originales de Estados ---
# df_gobierno = pd.DataFrame({
#     "NOM_EDO_PROD": ["Jalisco", "Jalisco", "Colima", "Colima", "Michoacán"],
#     "ACTUALIZADO": ["Si", "No", "Si", "No", "No"],
#     "Personas": [1500, 300, 800, 400, 1200]
# })

# fig1 = graficar_barras_porcentaje(
#     df=df_gobierno, 
#     col_x="NOM_EDO_PROD", 
#     col_color="ACTUALIZADO", 
#     col_valores="Personas",
#     titulo="Avance por Entidad Federativa"
# )
# st.plotly_chart(fig1, width='stretch')


# # --- EJEMPLO 2: Con datos de otra categoría completamente distinta ---
# df_ventas = pd.DataFrame({
#     "Sucursal": ["Norte", "Norte", "Sur", "Sur", "Este", "Este"],
#     "Tipo_Cliente": ["Nuevo", "Frecuente", "Nuevo", "Raro", "Nuevo", "Frecuente"],
#     "Personas": [200, 500, 300, 100, 400, 400]
# })

# fig2 = graficar_barras_porcentaje(
#     df=df_ventas, 
#     col_x="Sucursal", 
#     col_color="Tipo_Cliente", 
#     col_valores="Personas",
#     titulo="Distribución de Clientes por Sucursal"
# )
# st.plotly_chart(fig2, use_container_width=True)


# import streamlit as st
# import pandas as pd
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.dml.color import RGBColor
# from pptx.enum.text import PP_ALIGN
# from io import BytesIO

# # 1. Datos de ejemplo
# df = pd.DataFrame({
#     'Métrica': ['Usuarios Activos', 'Tasa de Conversión', 'Ingresos Totales', 'Churn Rate'],
#     'Este Mes': ['15,420', '3.4%', '$45,200', '2.1%'],
#     'Mes Anterior': ['14,100', '3.1%', '$41,000', '2.5%'],
#     'Cambio %': ['+9.3%', '+9.6%', '+10.2%', '-16.0%']
# })

# st.title("Exportar Tablas Nativas a PowerPoint")

# # Mostrar la tabla interactiva en Streamlit primero
# st.dataframe(df, use_container_width=True, hide_index=True)

# # 3. Función para construir el PPTX con Tablas Nativas y editables
# def crear_powerpoint(data_frame):
#     prs = Presentation()
#     # Usar una diapositiva en blanco (layout index 6 suele ser en blanco)
#     blank_slide_layout = prs.slide_layouts[6]
#     slide = prs.slides.add_slide(blank_slide_layout)
    
#     # Dimensiones de la tabla en la diapositiva
#     rows = len(data_frame) + 1  # Datos + Cabecera
#     cols = len(data_frame.columns)
    
#     left = Inches(0.5)
#     top = Inches(1.5)
#     width = Inches(9.0)
#     height = Inches(0.4 * rows)
    
#     # Agregar el objeto tabla nativo de PowerPoint
#     table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
#     table = table_shape.table
    
#     # Paleta de colores (Estilo Flextable / Corporativo)
#     azul_oscuro = RGBColor(30, 58, 138)
#     gris_claro = RGBColor(249, 250, 251)
#     blanco = RGBColor(255, 255, 255)
#     negro = RGBColor(0, 0, 0)
    
#     # --- Configurar Cabecera ---
#     for col_idx, col_name in enumerate(data_frame.columns):
#         cell = table.cell(0, col_idx)
#         cell.text = str(col_name)
#         cell.fill.solid()
#         cell.fill.fore_color.rgb = azul_oscuro
        
#         # Estilo del texto en la cabecera
#         for paragraph in cell.text_frame.paragraphs:
#             paragraph.alignment = PP_ALIGN.CENTER
#             for run in paragraph.runs:
#                 run.font.bold = True
#                 run.font.color.rgb = blanco
#                 run.font.name = 'Arial'
#                 run.font.size = Pt(13)

#     # --- Configurar Celdas de Datos (Editables) ---
#     for row_idx, row in data_frame.iterrows():
#         # Determinar color de fondo intercalado (zebra striping)
#         color_fondo = gris_claro if row_idx % 2 == 1 else blanco
        
#         for col_idx, value in enumerate(row):
#             cell = table.cell(row_idx + 1, col_idx)
#             cell.text = str(value)
#             cell.fill.solid()
#             cell.fill.fore_color.rgb = color_fondo
            
#             # Estilo del texto en las celdas
#             for paragraph in cell.text_frame.paragraphs:
#                 # Alinear a la izquierda la primera columna, a la derecha el resto
#                 paragraph.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.RIGHT
#                 for run in paragraph.runs:
#                     run.font.color.rgb = negro
#                     run.font.name = 'Arial'
#                     run.font.size = Pt(11)
                    
#     # Guardar la presentación en un buffer de memoria
#     binary_output = BytesIO()
#     prs.save(binary_output)
#     binary_output.seek(0)
#     return binary_output

# # 4. Botón de descarga en Streamlit
# ppt_data = crear_powerpoint(df)

# st.download_button(
#     label="📥 Descargar Reporte en PowerPoint (.pptx)",
#     data=ppt_data,
#     file_name="reporte_metricas.pptx",
#     mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
# )

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 1. Definición completa de los datos
data = [
    [9, "Ciudad de México", "4,203", "4,348", 3, "2,425", "2,490", "57.7%", "57.3%", "1,778"],
    [10, "Durango", "39,576", "54,298", 16, "34,500", "47,069", "87.2%", "86.7%", "5,076"],
    [17, "Morelos", "15,249", "18,322", 6, "11,987", "14,429", "78.6%", "78.8%", "3,262"],
    [18, "Nayarit", "29,261", "34,126", 9, "23,898", "27,265", "81.7%", "79.9%", "5,363"],
    [21, "Puebla", "121,230", "175,240", 23, "95,616", "133,581", "78.9%", "76.2%", "25,614"],
    [22, "Querétaro", "16,915", "24,041", 13, "14,065", "19,545", "83.2%", "81.3%", "2,850"],
    [29, "Tlaxcala", "25,700", "45,950", 10, "22,263", "38,389", "86.6%", "83.5%", "3,437"],
    [32, "Zacatecas", "68,413", "112,404", 15, "60,057", "100,083", "87.8%", "89.0%", "8,356"]
]

headers = [
    "Cve", "Representación", "Meta\npersonas", "Meta\npredios", 
    "No. de\nventanillas", "Avance\nregistro de\npersonas", 
    "Avance\nregistro de\npredios", "% de\navance\npersonas", 
    "% de\navance\npredios", "Personas\nNo\nactualizadas"
]

total_row = ["Total", "", "320,547", "468,729", "95", "264,811", "382,851", "82.6%", "81.7%", "55,736"]

# 2. Paleta cromática por bloques de columnas (Extraída de la imagen)
# Formato: (Color Cabecera, Color Fila Datos, Color Fila Totales)
COLOR_COLUMNAS = [
    (RGBColor(0, 50, 40), RGBColor(255, 255, 255), RGBColor(0, 40, 30)),      # Cve (Verde muy oscuro)
    (RGBColor(0, 50, 40), RGBColor(255, 255, 255), RGBColor(0, 40, 30)),      # Representación
    (RGBColor(0, 75, 60), RGBColor(245, 248, 245), RGBColor(0, 60, 45)),      # Meta personas (Verde Medio)
    (RGBColor(10, 65, 50), RGBColor(245, 248, 245), RGBColor(0, 50, 38)),     # Meta predios
    (RGBColor(115, 75, 40), RGBColor(253, 250, 245), RGBColor(90, 55, 25)),   # No. de ventanillas (Café)
    (RGBColor(50, 100, 80), RGBColor(240, 247, 243), RGBColor(40, 85, 65)),   # Avance personas (Verde claro)
    (RGBColor(45, 95, 75), RGBColor(240, 247, 243), RGBColor(35, 80, 60)),    # Avance predios
    (RGBColor(65, 115, 90), RGBColor(242, 249, 245), RGBColor(50, 95, 75)),   # % avance personas
    (RGBColor(55, 105, 80), RGBColor(242, 249, 245), RGBColor(42, 85, 65)),   # % avance predios
    (RGBColor(110, 70, 35), RGBColor(252, 248, 242), RGBColor(85, 50, 20))    # No actualizadas (Marrón/Ocre)
]

TEXTO_BLANCO = RGBColor(255, 255, 255)
TEXTO_NEGRO = RGBColor(40, 40, 40)

# 3. Inicializar presentación (Formato 16:9)
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts)

# Crear estructura de la tabla
rows, cols = len(data) + 2, len(headers)
left, top, width, height = Inches(0.4), Inches(1.5), Inches(12.53), Inches(4.5)
table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# 4. Rellenar y dar formato dinámico por celda utilizando la paleta
# Cabeceras (Fila 0)
for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_COLUMNAS[col_idx]  # Asigna color de cabecera de esa columna
    p = cell.text_frame.paragraphs
    p.text = header
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = TEXTO_BLANCO

# Filas de datos (Filas intermedias)
for row_idx, row_data in enumerate(data, start=1):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_COLUMNAS[col_idx]  # Asigna color de celda de esa columna
        p = cell.text_frame.paragraphs
        p.text = str(val)
        p.font.size = Pt(11)
        p.font.color.rgb = TEXTO_NEGRO
        p.alignment = PP_ALIGN.LEFT if col_idx == 1 else (PP_ALIGN.CENTER if col_idx == 0 or col_idx == 4 else PP_ALIGN.RIGHT)

# Fila de Totales (Última fila)
for col_idx, val in enumerate(total_row):
    cell = table.cell(rows - 1, col_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_COLUMNAS[col_idx]  # Asigna color de total de esa columna
    p = cell.text_frame.paragraphs
    p.text = str(val)
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = TEXTO_BLANCO
    p.alignment = PP_ALIGN.LEFT if col_idx == 0 else (PP_ALIGN.CENTER if col_idx == 4 else PP_ALIGN.RIGHT)

# 5. Ajustar anchos específicos de columnas para que encajen de forma idéntica
table.columns.width = Inches(0.5)   # Cve
table.columns.width = Inches(1.7)   # Representación
table.columns.width = Inches(1.1)   # Meta personas
table.columns.width = Inches(1.1)   # Meta predios
table.columns.width = Inches(1.0)   # Ventanillas
table.columns.width = Inches(1.3)   # Avance personas
table.columns.width = Inches(1.3)   # Avance predios
table.columns.width = Inches(1.1)   # % personas
table.columns.width = Inches(1.1)   # % predios
table.columns.width = Inches(1.3)   # No actualizadas

prs.save("Reporte_Avance_Color_Columnas.pptx")
print("Presentación con colores por columna creada exitosamente.")




################################################
# import streamlit as st
# import pandas as pd
# import numpy as np
# import plotly.graph_objects as go

# # 1. Configuración de la página de Streamlit
# st.set_page_config(page_title="Pirámide Poblacional", layout="centered")
# st.title("📊 Pirámide Poblacional Interactiva")
# st.write("Visualización demográfica construida con Streamlit y Plotly.")

# # 2. Creación de un dataset ficticio para demostración
# @st.cache_data
# def cargar_datos():
#     rangos_edad = [f"{i}-{i+4}" for i in range(0, 80, 5)] + ["80+"]
#     np.random.seed(42)
    
#     # Simular datos decrecientes para simular una pirámide estándar
#     base_hombres = np.linspace(500, 50, len(rangos_edad)) + np.random.randint(-30, 30, len(rangos_edad))
#     base_mujeres = np.linspace(510, 60, len(rangos_edad)) + np.random.randint(-30, 30, len(rangos_edad))
    
#     df = pd.DataFrame({
#         "Rango_Edad": rangos_edad,
#         "Hombres": base_hombres.astype(int),
#         "Mujeres": base_mujeres.astype(int)
#     })
#     return df

# df = cargar_datos()

# # Interactividad con widgets de Streamlit: Multiplicador simulador
# filtro_escala = st.slider("Simular crecimiento / escala de población:", 0.5, 2.0, 1.0, 0.1)
# df_filtrado = df.copy()
# df_filtrado["Hombres"] = (df_filtrado["Hombres"] * filtro_escala).astype(int)
# df_filtrado["Mujeres"] = (df_filtrado["Mujeres"] * filtro_escala).astype(int)

# # 3. Preparación de los datos para la pirámide (Valores negativos para hombres)
# valores_hombres = df_filtrado["Hombres"] * -1
# valores_mujeres = df_filtrado["Mujeres"]
# categorias_edad = df_filtrado["Rango_Edad"]

# # 4. Construcción del gráfico con Plotly Graph Objects
# fig = go.Figure()

# # Barra para Hombres (Izquierda)
# fig.add_trace(go.Bar(
#     y=categorias_edad,
#     x=valores_hombres,
#     name="Hombres",
#     orientation="h",
#     marker=dict(color="#1f77b4"),
#     hoverinfo="text",
#     text=df_filtrado["Hombres"].apply(lambda x: f"Hombres: {x:,}")
# ))

# # Barra para Mujeres (Derecha)
# fig.add_trace(go.Bar(
#     y=categorias_edad,
#     x=valores_mujeres,
#     name="Mujeres",
#     orientation="h",
#     marker=dict(color="#e377c2"),
#     hoverinfo="text",
#     text=df_filtrado["Mujeres"].apply(lambda x: f"Mujeres: {x:,}")
# ))

# # 5. Ajustes estéticos y de formato del eje X
# fig.update_layout(
#     barmode="overlay",
#     bargap=0.1,
#     title_text="Distribución por Edad y Sexo",
#     title_x=0.5,
#     xaxis=dict(
#         tickvals=[-1000, -750, -500, -250, 0, 250, 500, 750, 1000],
#         # Transformar las etiquetas negativas a positivas en la vista del usuario
#         ticktext=["1k", "750", "500", "250", "0", "250", "500", "750", "1k"],
#         title="Población"
#     ),
#     yaxis=dict(title="Rango de Edad"),
#     legend=dict(x=0.85, y=1)
# )

# # 6. Desplegar el gráfico en el dashboard de Streamlit
# st.plotly_chart(fig, use_container_width=True)

# # Mostrar la tabla de datos original si se desea inspeccionar
# if st.checkbox("Mostrar tabla de datos fuente"):
#     st.dataframe(df_filtrado, use_container_width=True)


# import pandas as pd
# import streamlit as st

# # 1. Creamos los datos demográficos de ejemplo
# df = pd.DataFrame({
#     "Rango Edad": ["18-25", "26-35", "36-45", "46+"],
#     "Usuarios": [1250, 3400, 2100, 850],
#     "Porcentaje": [0.16, 0.45, 0.28, 0.11]
# })

# # 2. Definimos las funciones de estilo (como los formateadores de flextable)
# def estilo_tabla(styler):
#     styler.set_table_styles([
#         # Estilo para los Headers (Fondo oscuro, texto grande y bold)
#         {"selector": "th", "props": [
#             ("background-color", "#1E3A8A"), 
#             ("color", "white"), 
#             ("font-size", "18px"), 
#             ("font-weight", "bold"),
#             ("text-align", "center"),
#             ("padding", "12px")
#         ]},
#         # Estilo para las Celdas (Letra grande, padding limpio)
#         {"selector": "td", "props": [
#             ("font-size", "16px"), 
#             ("padding", "10px"),
#             ("text-align", "center"),
#             ("border-bottom", "1px solid #E5E7EB")
#         ]},
#     ])
    
#     # Formatear números (Equivalente a colformat_num o colformat_pct)
#     styler.format({"Usuarios": "{:,}", "Porcentaje": "{:.1%}"})
    
#     # Color condicional (Resalta la fila con más usuarios como en flextable)
#     styler.highlight_max(subset=["Usuarios"], color="#D1FAE5")
    
#     return styler

# # 3. Aplicamos el estilo al DataFrame
# df_estilizado = df.style.pipe(estilo_tabla)

# # 4. Renderizamos en Streamlit convirtiendo a HTML
# st.subheader("📊 Perfil Demográfico Personalizado")

# # Convertimos el objeto Styler a HTML plano para que Streamlit mantenga el diseño CSS
# st.markdown(df_estilizado.to_html(index=False), unsafe_allow_html=True)




## Graficas con diseño mejorado

# import plotly.graph_objects as go

# # 1. Definir los datos
# header_values = ["Product Name", "Category", "Price", "Stock Status"]
# cell_values = [
#     ["Eco Water Bottle", "Home & Kitchen", "Premium Espresso", "Wireless Mouse"],
#     ["Sustainability", "Appliances", "Beverages", "Electronics"],
#     ["$25.00", "$120.00", "$45.00", "$15.00"],
#     ["In Stock", "Low Stock", "In Stock", "Out of Stock"]
# ]

# # Número total de filas de datos (en este ejemplo son 4)
# num_rows = len(cell_values[0])

# # 2. Crear la matriz de colores alternos para las filas
# # Fila 1: sin fondo (blanco/transparent) -> "white"
# # Fila 2: color personalizado -> "#D4C19C"
# row_colors = ["white", "#f8f4ed"] * (num_rows // 2 + 1)

# # 3. Crear el objeto de la tabla
# fig = go.Figure(data=[go.Table(
#     # columnwidth=[150, 150, 100, 120],
    
#     # Configuración del encabezado (mismo estilo anterior)
#     header=dict(
#         values=header_values,
#         fill_color="#235B4E",          # Fondo verde
#         font=dict(color="white", size=14, family="Arial"), # Texto blanco
#         line_color="#C29E5C",         # Borde dorado
#         align="center",
#         height=35
#     ),
    
#     # Configuración de las celdas con filas alternas
#     cells=dict(
#         values=cell_values,
#         fill_color=[row_colors],       # Se aplica la lista de colores alternos
#         font=dict(color="#333333", size=12, family="Arial"),
#         line_color="white",         # Borde dorado
#         align="left",
#         height=30
#     )
# )])

# # 4. Ajustar márgenes
# fig.update_layout(
#     margin=dict(l=10, r=10, t=10, b=10),
#     height=400
# )

# fig.show()

# import polars as pl
# import pandas as pd
# import matplotlib.pyplot as plt
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.dml.color import RGBColor

# # 1. CARGAR DATOS DESDE PARQUET
# df_raw = pl.read_parquet("concentrado_actualizados.parquet").to_pandas()
# df_raw.columns = df_raw.columns.str.lower()  # Forzar minúsculas por seguridad

# # Limpieza de nulos inicial
# df_raw = df_raw.dropna(subset=['mes', 'personas'])

# # ==========================================
# # ¡SOLUCIÓN!: AGRUPAR DATOS POR MES
# # ==========================================
# # Agrupamos por 'mes' y sumamos la columna 'personas'. 
# # .reset_index() mantiene 'mes' como una columna normal de la tabla.
# df = df_raw.groupby('mes', as_index=False)['personas'].sum()

# # Opcional: Si tus meses son texto o números, puedes ordenarlos aquí para que el gráfico sea coherente
# # df = df.sort_values(by='mes') 

# print("Datos agrupados listos para graficar y tabular:")
# print(df)

# # 2. GENERAR Y GUARDAR EL GRÁFICO CON MATPLOTLIB
# plt.figure(figsize=(7, 4))
# plt.plot(df['mes'], df['personas'], marker='o', color='#235B4E', linewidth=2.5)
# plt.title("Reporte de Ventas Mensuales", fontsize=14, pad=15)
# plt.grid(True, linestyle='--', alpha=0.5)
# plt.tight_layout()

# print("Generando gráfico con Matplotlib...")
# grafico_path = "grafico_matplotlib.png"
# plt.savefig(grafico_path, dpi=150)
# plt.close()

# # 3. CONFIGURAR PRESENTACIÓN Y COLORES
# prs = Presentation("plantilla 2026.pptx")
# blank_layout = prs.slide_layouts[5] 

# COLOR_VERDE = RGBColor(35, 91, 78)     # #235B4E
# COLOR_BEIGE = RGBColor(212, 193, 156)  # #D4C19C
# COLOR_BLANCO = RGBColor(255, 255, 255)

# # --- DIAPOSITIVA 1: INSERTAR GRÁFICO ---
# slide_grafico = prs.slides.add_slide(blank_layout)
# slide_grafico.shapes.add_picture(grafico_path, Inches(1.5), Inches(1.5), width=Inches(7))
# print("Diapositiva 1 (Gráfico) añadida con éxito.")

# # --- DIAPOSITIVA 2: CREAR TABLA ---
# slide_tabla = prs.slides.add_slide(blank_layout)
# rows, cols = len(df) + 1, len(df.columns)

# # Ahora 'rows' tiene una longitud segura determinada por la cantidad de meses únicos
# table = slide_tabla.shapes.add_table(rows, cols, Inches(1.5), Inches(1.5), Inches(7), Inches(0.4 * rows)).table
# print(f"Creando tabla resumida de {rows} filas por {cols} columnas...")

# # Llenar Encabezados
# for col_idx, col_name in enumerate(df.columns):
#     cell = table.cell(0, col_idx)
#     cell.text = str(col_name).upper()
#     cell.fill.solid()
#     cell.fill.fore_color.rgb = COLOR_VERDE
#     cell.text_frame.paragraphs[0].font.color.rgb = COLOR_BLANCO
#     cell.text_frame.paragraphs[0].font.bold = True

# # Llenar Datos con Filas Alternas
# for row_idx, row_data in df.iterrows():
#     for col_idx, value in enumerate(row_data):
#         cell = table.cell(row_idx + 1, col_idx)
        
#         # Formato numérico: si es la columna de personas (números), aplicamos formato entero con comas
#         if isinstance(value, (int, float)) and df.columns[col_idx] == 'personas':
#             cell.text = f"{value:,.0f}"  # Cambiado a enteros sin signo de $, solo comas de miles
#         else:
#             cell.text = str(value)
            
#         cell.fill.solid()
#         cell.fill.fore_color.rgb = COLOR_BEIGE if (row_idx + 1) % 2 == 0 else COLOR_BLANCO
#         cell.text_frame.paragraphs[0].font.size = Pt(11)

# # Guardar archivo final de forma segura
# prs.save("presentacion_final.pptx")
# print("¡Presentación 'presentacion_final.pptx' creada con éxito!")

import matplotlib.pyplot as plt

def crear_dona(df_dona, titulo, colores_lista=None, height=480):
    # 1. Configuración básica de estilos
    colores = colores_lista[:len(df_dona)] if colores_lista else plt.cm.tab10.colors
    font = FONT_FAMILY if 'FONT_FAMILY' in globals() else 'Noto Sans'
    color_txt = GUINDA if 'GUINDA' in globals() else '#6A1B29'
    
    # 2. Inicializar figura limpia (sin ejes traseros)
    fig, ax = plt.subplots(figsize=(6, height / 100), facecolor='none')
    
    # 3. Dibujar la dona y los porcentajes internos
    wedges, _, autotexts = ax.pie(
        df_dona['Personas'], 
        autopct='%1.1f%%', pctdistance=0.75, colors=colores, startangle=90,
        wedgeprops=dict(width=0.4, edgecolor='white', linewidth=2) # Anillo estilizado
    )
    
    # Estilo rápido para los porcentajes internos
    plt.setp(autotexts, color='white', fontfamily=font, weight='bold', size=11)

    # 4. Texto central (Total)
    ax.text(0, 0, f"{df_dona['Personas'].sum():,.0f}", ha='center', va='center', 
            fontsize=20, color=color_txt, fontfamily=font, weight='bold')

    # 5. Título y Leyenda simplificada abajo
    ax.set_title(titulo, fontsize=16, color=color_txt, fontfamily=font, pad=15)
    
    ax.legend(
        wedges, df_dona['Categoria'], 
        loc='upper center', bbox_to_anchor=(0.5, -0.05), 
        ncol=len(df_dona), frameon=False, prop={'family': font, 'size': 10}
    )
    ax.axis('equal')
    return fig



from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Inicializar la presentación
prs = Presentation("plantilla 2026.pptx")

# --- DIAPOSITIVA 1: TÍTULO ---
# Usar el diseño de portada predeterminado (índice 0)
slide1 = prs.slides.add_slide(prs.slide_layouts[0])

# 🚀 FUNCIÓN MÁGICA: Pone texto en coordenadas con formato personalizado
# Ahora agregamos "ancho=5" al final de la primera línea (por defecto será de 5 pulgadas)
def agregar_texto(slide, x=1, y=1, texto="", fuente="Noto Sans", tamano=14, negrita=True, color=(0,0,0), ancho=10):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(ancho), Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True # IMPORTANTE: Si el texto supera el 'ancho', saltará de línea
    p = tf.paragraphs[0]
    p.text = texto
    p.font.name = fuente
    p.font.size = Pt(tamano)
    p.font.bold = negrita
    p.font.color.rgb = RGBColor(*color) if isinstance(color, tuple) else RGBColor(0, 0, 0)  # Soporta tuplas RGB o nombres de color
    
    # ... (aquí va el resto de tu código del procesador de color) ...

# 📍 EJEMPLOS DE USO EN UNA SOLA LÍNEA
# Parámetros: (slide, X, Y, "Texto", "Fuente", Tamaño, Negrita, (R,G,B))

# Asignar textos
slide1.shapes.title.text = "Producción para el Bienestar"
slide1.placeholders[1].text = "Actualización o Corroboración de Datos e Integración de Expedientes - 2026"
agregar_texto(slide1, 10.5, 6.5, f"Fecha: {datetime.now().strftime('%d-%m-%Y')}",tamano=16, color = (255, 255, 255), ancho=2.5)   # Rojo

# --- DIAPOSITIVA 2: CONTENIDO ---
# Usar diseño de título y cuerpo (índice 1)
slide_layout_content = prs.slide_layouts[3]
slide2 = prs.slides.add_slide(slide_layout_content)

# Título de la segunda diapositiva
slide2.shapes.title.text = "Resultados del Trimestre"

# Añadir un cuadro de texto personalizado
left = Inches(1)
top = Inches(2.5)
width = Inches(8)
height = Inches(4)

txBox = slide2.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "Este es el párrafo principal del reporte."

# Añadir un párrafo secundario con formato
p = tf.add_paragraph()
p.text = "• Incremento del 15% en ventas masivas."
p.font.bold = True
p.font.size = Pt(18)



# 1. Definir los datos de la tabla (Encabezados + Contenido + Totales)
datos = [
    ["Cve", "Representación", "Meta personas", "Meta predios", "No. de ventanillas", "Avance registro de personas", "Avance registro de predios", "% de avance personas", "% de avance predios", "Personas No actualizadas"],
    ["9", "Ciudad de México", "4,203", "4,348", "3", "2,425", "2,490", "57.7%", "57.3%", "1,778"],
    ["10", "Durango", "39,576", "54,298", "16", "34,500", "47,069", "87.2%", "86.7%", "5,076"],
    ["17", "Morelos", "15,249", "18,322", "6", "11,987", "14,429", "78.6%", "78.8%", "3,262"],
    ["18", "Nayarit", "29,261", "34,126", "9", "23,888", "27,265", "81.7%", "79.9%", "5,363"],
    ["21", "Puebla", "121,230", "175,240", "23", "95,616", "133,581", "78.9%", "76.2%", "25,614"],
    ["22", "Querétaro", "16,915", "24,041", "13", "14,065", "19,545", "83.2%", "81.3%", "2,850"],
    ["29", "Tlaxcala", "25,700", "45,950", "10", "22,263", "38,389", "86.6%", "83.5%", "3,437"],
    ["32", "Zacatecas", "68,413", "112,404", "15", "60,057", "100,083", "87.8%", "89.0%", "8,356"],
    ["Total", "", "320,547", "468,729", "95", "264,811", "382,851", "82.6%", "81.7%", "55,736"]
]

# 2. Configurar dimensiones y posición de la tabla
filas = len(datos)
columnas = len(datos[0])

x = Inches(0.5)      # Posición izquierda
y = Inches(1.5)      # Posición superior
ancho = Inches(12)   # Ancho total de la tabla (ajustado para pantallas panorámicas 16:9)
alto = Inches(4.5)   # Alto total estimado de la tabla

# 3. Insertar la tabla en la diapositiva
shape = slide2.shapes.add_table(filas, columnas, x, y, ancho, alto)
table = shape.table

# 4. Ajustar anchos específicos de columnas (Opcional, para que se vea ordenado)
table.columns[0].width = Inches(0.5)  # Columna Cve más angosta
table.columns[1].width = Inches(2.0)  # Columna Representación más ancha

# 1. Datos exactos de tu imagen
datos = [
    ["Cve", "Representación", "Meta personas", "Meta predios", "No. de ventanillas", "Avance registro de personas", "Avance registro de predios", "% de avance personas", "% de avance predios", "Personas No actualizadas"],
    ["9", "Ciudad de México", "4,203", "4,348", "3", "2,425", "2,490", "57.7%", "57.3%", "1,778"],
    ["10", "Durango", "39,576", "54,298", "16", "34,500", "47,069", "87.2%", "86.7%", "5,076"],
    ["17", "Morelos", "15,249", "18,322", "6", "11,987", "14,429", "78.6%", "78.8%", "3,262"],
    ["18", "Nayarit", "29,261", "34,126", "9", "23,898", "27,265", "81.7%", "79.9%", "5,363"],
    ["21", "Puebla", "121,230", "175,240", "23", "95,616", "133,581", "78.9%", "76.2%", "25,614"],
    ["22", "Querétaro", "16,915", "24,041", "13", "14,065", "19,545", "83.2%", "81.3%", "2,850"],
    ["29", "Tlaxcala", "25,700", "45,950", "10", "22,263", "38,389", "86.6%", "83.5%", "3,437"],
    ["32", "Zacatecas", "68,413", "112,404", "15", "60,057", "100,083", "87.8%", "89.0%", "8,356"],
    ["Total", "", "320,547", "468,729", "95", "264,811", "382,851", "82.6%", "81.7%", "55,736"]
]

# 2. Crear tabla (10 filas x 10 columnas) en coordenadas (X=0.5, Y=1.5, Ancho=12, Alto=4.5)
tabla = slide2.shapes.add_table(10, 10, Inches(0.5), Inches(1.5), Inches(12), Inches(4.5)).table

# 3. Rellenar de forma masiva y compacta
for f_idx, fila in enumerate(datos):
    for c_idx, valor in enumerate(fila):
        tabla.cell(f_idx, c_idx).text = valor

# 🚀 FUNCIÓN MÁGICA: Crea una fila de tarjetas con encabezado y viñetas
def crear_bloques_horizontales(slide, x_ini, y_ini, datos, color_header=(166, 30, 77), ancho_bloque=2.8):
    alto_tit = 0.6
    alto_body = 2.5
    separacion = 0.3
    
    for i, (titulo, vinetas) in enumerate(datos):
        # 📐 Calcular posición horizontal dinámica
        x = Inches(x_ini + i * (ancho_bloque + separacion))
        
        # 🟥 1. Crear el Encabezado
        b_tit = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(y_ini), Inches(ancho_bloque), Inches(alto_tit))
        b_tit.fill.solid()
        b_tit.fill.fore_color.rgb = RGBColor(*color_header)
        b_tit.line.fill.background() # Quita el borde
        b_tit.text_frame.text = titulo
        b_tit.text_frame.paragraphs[0].font.size = Pt(12)
        b_tit.text_frame.paragraphs[0].font.bold = True
        b_tit.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # ⬜ 2. Crear el Cuerpo Gris
        b_body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(y_ini + alto_tit), Inches(ancho_bloque), Inches(alto_body))
        b_body.fill.solid()
        b_body.fill.fore_color.rgb = RGBColor(240, 242, 245)
        b_body.line.color.rgb = RGBColor(200, 200, 200) # Borde gris claro
        
        tf = b_body.text_frame
        tf.word_wrap = True
        
        # 📑 3. Inyectar las viñetas
        for j, texto in enumerate(vinetas):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"• {texto}"
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(30, 30, 30)

# ==========================================================
# 📍 CONFIGURACIÓN DE TUS DATOS (Estructura limpia)
# ==========================================================
informacion = [
    ["Constancia de Posesión", ["7% de los productores actualizados, presentaron constancias"]],
    ["RENTA", ["El 13.6% presentó documento de RENTA.", "Ahora se conoce la fecha."]],
    ["Documentos", ["0.4% presentó documentos no válidos.", "99.6% presentó papeles aceptables."]],
    ["Género", ["El 35.2% son mujeres.", "El 64.8% son hombres."]]
]

# 🎯 INVOCACIÓN EN UNA SOLA LÍNEA
# Parámetros: (slide, X_inicial, Y_inicial, datos, color_rgb_opcional)
crear_bloques_horizontales(slide1, 0.5, 2.0, informacion)




# 3. Guardar el archivo
prs.save("presentacion_python.pptx")
print("¡Presentación creada con éxito!")


import io
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches

# 1. Crear una figura muy simple con Plotly
fig = go.Figure(
    data=[go.Bar(x=['A', 'B', 'C'], y=[10, 20, 15], marker_color='#6A1B29')],
    layout=go.Layout(
        title="Gráfico de Ejemplo",
        width=500,
        height=400,
        paper_bgcolor="rgba(0,0,0,0)", # Fondo transparente
        plot_bgcolor="rgba(0,0,0,0)"
    )
)

# 2. Exportar el gráfico a un buffer de memoria usando Kaleido
image_stream = io.BytesIO()
fig.write_image(image_stream, format="png", scale=2)
image_stream.seek(0)

# 3. Configurar PowerPoint e insertar la imagen en la diapositiva 1
prs = Presentation()
slide1 = prs.slides.add_slide(prs.slide_layouts[0]) # Usando tu diseño slide_layouts[0]

# Posición en la diapositiva (en pulgadas)
left = Inches(1.5)
top = Inches(2.5)
width = Inches(7.0)

# Insertar la imagen directamente desde la memoria
slide1.shapes.add_picture(image_stream, left, top, width=width)

# 4. Guardar la presentación
prs.save("presentacion_plotly.pptx")
print("¡Presentación creada con éxito junto con el gráfico de Plotly!")

import io
import plotly.graph_objects as go
import plotly.io as pio
from pptx import Presentation
from pptx.util import Inches

# 1. Supongamos que tienes una lista de diccionarios con tus datos
datos_reporte = [
    {"titulo": "Ventas Región Norte", "categorias": ['A', 'B', 'C'], "valores": [10, 20, 15]},
    {"titulo": "Ventas Región Sur", "categorias": ['D', 'E', 'F'], "valores": [30, 12, 45]},
    {"titulo": "Ventas Región Centro", "categorias": ['G', 'H', 'I'], "valores": [25, 18, 22]}
]

# 2. Inicializar la presentación una sola vez antes del bucle
prs = Presentation()
diseño_diapositiva = prs.slide_layouts[5] # Diseño típico con título arriba y espacio blanco abajo

# 3. Bucle para generar los gráficos e insertarlos
for data in datos_reporte:
    # Crear la figura de Plotly para esta iteración
    fig = go.Figure(
        data=[go.Bar(x=data["categorias"], y=data["valores"], marker_color='#6A1B29')],
        layout=go.Layout(
            title=data["titulo"],
            width=600, height=450,
            paper_bgcolor="rgba(0,0,0,0)",
            plot_bgcolor="rgba(0,0,0,0)"
        )
    )
    
    # Agregar una nueva diapositiva para este gráfico
    nueva_slide = prs.slides.add_slide(diseño_diapositiva)
    
    # Exportar a memoria e insertar (El 'with' libera la memoria al terminar cada vuelta)
    with io.BytesIO() as image_stream:
        fig.write_image(image_stream, format="png", scale=2)
        image_stream.seek(0)
        
        # Insertar en la diapositiva actual
        nueva_slide.shapes.add_picture(
            image_stream, 
            left=Inches(1.5), 
            top=Inches(2.0), 
            width=Inches(7.0)
        )
    
    print(f"Diapositiva '{data['titulo']}' procesada con éxito.")

# 4. Guardar el archivo final al salir del bucle
prs.save("reporte_mensual_automatizado.pptx")

print("\n¡Presentación completa generada y memoria 100% liberada!")
