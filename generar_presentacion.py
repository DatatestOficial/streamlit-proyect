from datetime import datetime
from zoneinfo import ZoneInfo

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches, Cm
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

import streamlit as st
import psycopg 
import io
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from babel.dates import format_date



# Valores por defecto basados en la paleta de la imagen
RGB_VERDE = RGBColor(16,49,43)
RGB_VERDE_CLARO = RGBColor(35,91,78)
RGB_ROJO = RGBColor(105,28,50)
# RGB_ROJO_CLARO = RGBColor(159,34,65)
RGB_ROJO_CLARO = RGBColor(156,33,72)
RGB_CAFE = RGBColor(194,158,92)
RGB_GUINDA = RGB_ROJO   # Fondo de la cabecera
RGB_DORADO = RGBColor(190, 151, 91)   # Cifra central
RGB_BLANCO = RGBColor(255, 255, 255)   # Texto de la cabecera y fondo tarjeta
RGB_GRIS_BORDE = RGBColor(30, 41, 59) # Borde sutil inferior si se desea
RGB_GRIS = RGBColor(118, 113, 113) # Texto de la cabecera y fondo tarjeta

FONT_FAMILY = "Noto Sans"
FONT_SIZE_AXIS = 18
FONT_SIZE_TITLE = 24
PALETA_INSTITUCIONAL = ["#10312B", "#691C32", "#C29E5C", "#235B4E", "#9F2241", "#D4C19C", "#44546A", "#52492E", "#52492E", "#f8f4ed"]
GUINDA = "#621132"
GUINDA_CLARO = "#9F2241"
DORADO = "#D4C19C"
VERDE = "#285C4D"
AMARILLO = "#745526"
VERDE_CLARO = "#3A7D6B"
CREMA = "#F5F1EB"


# =====================================================
# Funciones para pptx
# =====================================================
# Funciones
@st.cache_data(ttl=3600*6)
def cargar_datos(query, parametros=None):
    with psycopg.connect(st.secrets["supabase"]["DATABASE_URL"]) as conn:
        with conn.cursor() as cur:
            cur.execute(query, parametros)
            filas = cur.fetchall()
            columnas = [desc.name for desc in cur.description]
    return pd.DataFrame(filas, columns=columnas).reset_index(drop=True)


def add_header_row(
    slide,
    left,
    top,
    width,
    text,
    header_colors,
    column_widths,
    column_alignments=None,
    row_height=0.65,
    font_name="Noto Sans",
    font_size=14,
    font_color=(255, 255, 255),
):
    """
    Crea una fila de encabezados institucionales.

    Parameters
    ----------
    slide : pptx.slide.Slide

    left, top, width :
        Posición y ancho total.

    text : list[str]
        Textos de encabezado.

    header_colors : list[str]
        Colores HEX por columna.

    column_widths : list[float]
        Proporciones de ancho.
        Ejemplo:
        [40,20,20,20]

    column_alignments : list[str]
        ["left","center","right"]

    row_height : float
        Alto en cm.

    font_name : str

    font_size : int

    font_color : tuple
        RGB
        Ejemplo:
        (255,255,255)
    """

    # ==========================================
    # VALIDACIONES
    # ==========================================

    ncols = len(text)

    if len(header_colors) != ncols:
        raise ValueError(
            "header_colors debe tener "
            "la misma longitud que text"
        )

    if len(column_widths) != ncols:
        raise ValueError(
            "column_widths debe tener "
            "la misma longitud que text"
        )

    if column_alignments is None:
        column_alignments = ["center"] * ncols

    if len(column_alignments) != ncols:
        raise ValueError(
            "column_alignments debe tener "
            "la misma longitud que text"
        )

    # ==========================================
    # BORDE
    # ==========================================

    def set_border(cell):

        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

        for edge in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:

            ln = OxmlElement(edge)
            ln.set("w", "12700")

            solid = OxmlElement("a:solidFill")

            srgb = OxmlElement("a:srgbClr")
            srgb.set("val", "C29E5C")

            solid.append(srgb)
            ln.append(solid)

            tcPr.append(ln)

    # ==========================================
    # ALINEACIÓN
    # ==========================================

    def get_alignment(value):

        value = value.lower()

        if value == "left":
            return PP_ALIGN.LEFT

        if value == "center":
            return PP_ALIGN.CENTER

        if value == "right":
            return PP_ALIGN.RIGHT

        raise ValueError(
            f"Alineación inválida: {value}"
        )

    # ==========================================
    # TABLA DE UNA SOLA FILA
    # ==========================================

    table = slide.shapes.add_table(
        1,
        ncols,
        left,
        top,
        width,
        Cm(row_height)
    ).table

    # ==========================================
    # ANCHOS
    # ==========================================

    total_prop = sum(column_widths)

    for i, p in enumerate(column_widths):

        table.columns[i].width = int(
            width * p / total_prop
        )

    table.rows[0].height = Cm(row_height)

    # ==========================================
    # HEADER
    # ==========================================

    for c in range(ncols):

        cell = table.cell(0, c)

        color = (
            header_colors[c]
            .replace("#", "")
        )

        cell.text = str(text[c])

        cell.fill.solid()

        cell.fill.fore_color.rgb = RGBColor(
            int(color[0:2], 16),
            int(color[2:4], 16),
            int(color[4:6], 16)
        )

        cell.vertical_anchor = (
            MSO_VERTICAL_ANCHOR.MIDDLE
        )

        cell.margin_left = Cm(0.02)
        cell.margin_right = Cm(0.02)

        tf = cell.text_frame

        for p in tf.paragraphs:
            p.alignment = get_alignment(
                column_alignments[c]
            )

        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(
            *font_color
        )

        set_border(cell)

    return table

def add_table_body(
    slide,
    df,
    left,
    top,
    width,
    height,
    font_size=14,
    column_widths=None,
    column_alignments=None,
    row_height=0.55,
    font_name="Noto Sans",
):
    """
    Tabla institucional SIN encabezado y SIN fila total.
    """

    BODY_TEXT = RGBColor(0, 0, 0)

    ZEBRA_FILL = RGBColor(248, 244, 237)

    BORDER_HEX = "C29E5C"

    # ==========================================
    # VALIDACIONES
    # ==========================================

    cols = len(df.columns)

    if column_widths is not None:

        if len(column_widths) != cols:

            raise ValueError(
                f"column_widths debe tener "
                f"{cols} elementos"
            )

    if column_alignments is not None:

        if len(column_alignments) != cols:

            raise ValueError(
                f"column_alignments debe tener "
                f"{cols} elementos"
            )

    # ==========================================
    # BORDE
    # ==========================================

    def set_border(cell):

        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

        for edge in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:

            ln = OxmlElement(edge)
            ln.set("w", "12700")

            solid = OxmlElement("a:solidFill")

            srgb = OxmlElement("a:srgbClr")
            srgb.set("val", BORDER_HEX)

            solid.append(srgb)

            ln.append(solid)

            tcPr.append(ln)

    # ==========================================
    # FORMATO
    # ==========================================

    def format_value(value):

        if pd.isna(value):
            return ""

        if isinstance(value, (int, float)):
            return f"{value:,.0f}"

        return str(value)

    # ==========================================
    # ALINEACIÓN
    # ==========================================

    def get_alignment(value):

        value = str(value).lower()

        if value == "left":
            return PP_ALIGN.LEFT

        if value == "center":
            return PP_ALIGN.CENTER

        if value == "right":
            return PP_ALIGN.RIGHT

        raise ValueError(
            f"Alineación inválida: {value}"
        )

    # ==========================================
    # TABLA
    # ==========================================

    rows = len(df)

    table = slide.shapes.add_table(
        rows,
        cols,
        left,
        top,
        width,
        height
    ).table

    # ==========================================
    # ANCHOS
    # ==========================================

    if column_widths is not None:

        total_prop = sum(column_widths)

        for i, p in enumerate(column_widths):

            table.columns[i].width = int(
                width * p / total_prop
            )

    # ==========================================
    # ALTURA
    # ==========================================

    for row in table.rows:
        row.height = Cm(row_height)

    numeric_cols = df.select_dtypes(
        include="number"
    ).columns

    # ==========================================
    # CUERPO
    # ==========================================

    for r in range(len(df)):

        for c, col in enumerate(df.columns):

            cell = table.cell(r, c)

            value = df.iloc[r, c]

            cell.text = format_value(value)

            cell.vertical_anchor = (
                MSO_VERTICAL_ANCHOR.MIDDLE
            )

            cell.margin_left = Cm(0.02)
            cell.margin_right = Cm(0.02)

            # Zebra
            if r % 2 == 0:

                cell.fill.solid()

                cell.fill.fore_color.rgb = (
                    ZEBRA_FILL
                )

            else:

                cell.fill.background()

            p = cell.text_frame.paragraphs[0]

            p.font.name = font_name
            p.font.size = Pt(font_size)
            p.font.bold = False
            p.font.color.rgb = BODY_TEXT

            # Alineación personalizada
            if column_alignments is not None:

                p.alignment = get_alignment(
                    column_alignments[c]
                )

            else:

                if col in numeric_cols:
                    p.alignment = PP_ALIGN.RIGHT
                else:
                    p.alignment = PP_ALIGN.LEFT

            set_border(cell)

    return table

def add_institutional_table(
    slide,
    df,
    left,
    top,
    width,
    height,
    font_size=14,
    header_colors=None,
    column_widths=None,
    column_alignments=None,
    row_height=0.55,
    font_name="Noto Sans",
    add_total=False,
):
    """
    Tabla institucional para PowerPoint.

    Parameters
    ----------
    slide : pptx.slide.Slide

    df : pandas.DataFrame

    left, top, width, height :
        Posición y tamaño de la tabla.

    font_size : int
        Tamaño de fuente para encabezado y cuerpo.

    header_colors : list[str]
        Lista de colores HEX para encabezados.
        Ejemplo:
        ["#10312B","#691C32","#C29E5C"]

        Si hay más columnas que colores,
        los colores se repiten.

    column_widths : list[float]
        Proporciones de ancho.
        Ejemplo:
        [40,20,20,20]

    column_alignments : list[str]
        Alineación por columna.

        Valores válidos:
        ["left","center","right"]

        Si es None:
        - numéricos -> right
        - texto -> left

    row_height : float
        Alto de fila en centímetros.

    font_name : str

    add_total : bool
        Agrega fila TOTAL.
    """

    # ==================================================
    # PALETA DEFAULT
    # ==================================================

    if header_colors is None:
        header_colors = [
            "#10312B",
            "#691C32",
            "#C29E5C",
            "#235B4E",
            "#9F2241",
            "#D4C19C",
            "#44546A",
            "#52492E",
            "#52492E",
            "#F8F4ED",
        ]

    # ==================================================
    # COLORES
    # ==================================================

    HEADER_TEXT = RGBColor(255, 255, 255)

    BODY_TEXT = RGBColor(0, 0, 0)

    BORDER_HEX = "C29E5C"

    ZEBRA_FILL = RGBColor(248, 244, 237)

    # ==================================================
    # VALIDACIONES
    # ==================================================

    cols_original = len(df.columns)

    if column_widths is not None:

        if len(column_widths) != cols_original:

            raise ValueError(
                f"column_widths debe tener "
                f"{cols_original} elementos"
            )

    if column_alignments is not None:

        if len(column_alignments) != cols_original:

            raise ValueError(
                f"column_alignments debe tener "
                f"{cols_original} elementos"
            )

    # ==================================================
    # BORDE
    # ==================================================

    def set_border(cell):

        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

        for edge in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:

            ln = OxmlElement(edge)
            ln.set("w", "12700")

            solid = OxmlElement("a:solidFill")

            srgb = OxmlElement("a:srgbClr")
            srgb.set("val", BORDER_HEX)

            solid.append(srgb)
            ln.append(solid)

            tcPr.append(ln)

    # ==================================================
    # FORMATO
    # ==================================================

    def format_value(value):

        if pd.isna(value):
            return ""

        if isinstance(value, (int, float)):
            return f"{value:,.0f}"

        return str(value)

    # ==================================================
    # ALINEACIÓN
    # ==================================================

    def get_alignment(value):

        value = str(value).lower()

        if value == "left":
            return PP_ALIGN.LEFT

        if value == "center":
            return PP_ALIGN.CENTER

        if value == "right":
            return PP_ALIGN.RIGHT

        raise ValueError(
            f"Alineación no válida: {value}"
        )

    # ==================================================
    # DATAFRAME
    # ==================================================

    work_df = df.copy()

    numeric_cols = work_df.select_dtypes(
        include="number"
    ).columns

    # ==================================================
    # FILA TOTAL
    # ==================================================

    if add_total:

        total_row = {}

        for col in work_df.columns:

            if col in numeric_cols:
                total_row[col] = work_df[col].sum()
            else:
                total_row[col] = ""

        total_row[work_df.columns[0]] = "TOTAL"

        work_df = pd.concat(
            [work_df, pd.DataFrame([total_row])],
            ignore_index=True
        )

    # ==================================================
    # CREAR TABLA
    # ==================================================

    rows = len(work_df) + 1
    cols = len(work_df.columns)

    table = slide.shapes.add_table(
        rows,
        cols,
        left,
        top,
        width,
        height
    ).table

    # ==================================================
    # ANCHOS
    # ==================================================

    if column_widths is not None:

        total_prop = sum(column_widths)

        for i, p in enumerate(column_widths):

            table.columns[i].width = int(
                width * p / total_prop
            )

    # ==================================================
    # ALTURA FILAS
    # ==================================================

    for row in table.rows:
        row.height = Cm(row_height)

    # ==================================================
    # HEADER
    # ==================================================

    for c, col_name in enumerate(work_df.columns):

        cell = table.cell(0, c)

        cell.text = str(col_name)

        color = header_colors[c % len(header_colors)]
        color = color.replace("#", "")

        cell.fill.solid()

        cell.fill.fore_color.rgb = RGBColor(
            int(color[0:2], 16),
            int(color[2:4], 16),
            int(color[4:6], 16)
        )

        cell.vertical_anchor = (
            MSO_VERTICAL_ANCHOR.MIDDLE
        )

        cell.margin_left = Cm(0.02)
        cell.margin_right = Cm(0.02)

        p = cell.text_frame.paragraphs[0]

        p.alignment = PP_ALIGN.CENTER

        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = HEADER_TEXT

        set_border(cell)

    # ==================================================
    # CUERPO
    # ==================================================

    for r in range(len(work_df)):

        is_total = (
            add_total
            and
            r == len(work_df) - 1
        )

        for c, col in enumerate(work_df.columns):

            value = work_df.iloc[r, c]

            cell = table.cell(r + 1, c)

            cell.text = format_value(value)

            cell.vertical_anchor = (
                MSO_VERTICAL_ANCHOR.MIDDLE
            )

            cell.margin_left = Cm(0.02)
            cell.margin_right = Cm(0.02)

            p = cell.text_frame.paragraphs[0]

            p.font.name = font_name
            p.font.size = Pt(font_size)

            # =====================================
            # FILA TOTAL
            # =====================================

            if is_total:

                p.alignment = PP_ALIGN.CENTER

                color = header_colors[
                    c % len(header_colors)
                ]

                color = color.replace("#", "")

                cell.fill.solid()

                cell.fill.fore_color.rgb = RGBColor(
                    int(color[0:2], 16),
                    int(color[2:4], 16),
                    int(color[4:6], 16)
                )

                p.font.bold = True
                p.font.color.rgb = HEADER_TEXT

            # =====================================
            # CUERPO
            # =====================================

            else:

                # Zebra institucional
                if r % 2 == 0:

                    cell.fill.solid()

                    cell.fill.fore_color.rgb = (
                        ZEBRA_FILL
                    )

                else:

                    cell.fill.background()

                p.font.bold = False
                p.font.color.rgb = BODY_TEXT

                # -------------------------
                # Alineación personalizada
                # -------------------------

                if column_alignments is not None:

                    p.alignment = get_alignment(
                        column_alignments[c]
                    )

                else:

                    if col in numeric_cols:
                        p.alignment = PP_ALIGN.RIGHT
                    else:
                        p.alignment = PP_ALIGN.LEFT

            set_border(cell)

    return table

def crear_barras_porcentaje(
    df: pd.DataFrame,
    col_x: str,
    col_color: str,
    col_valores: str = "Personas",
    titulo: str = "Avance",
    lista_colores: list = None,
    orden_ascendente: bool = True,
    invertir_apilado: bool = True,
    font_size = 22,
    height = 680,
) -> go.Figure:

    # ----------------------------
    # VALIDACIÓN SEGURA
    # ----------------------------
    if not isinstance(df, pd.DataFrame) or df.empty:
        fig = go.Figure()
        fig.update_layout(
            title=f"{titulo} (Sin datos disponibles)",
            paper_bgcolor="rgba(0,0,0,0)",
            plot_bgcolor="rgba(0,0,0,0)"
        )
        return fig

    # ----------------------------
    # AGRUPACIÓN BASE
    # ----------------------------
    df_agrupado = (
        df.groupby([col_x, col_color], as_index=False)[col_valores]
        .sum()
    )
    # ----------------------------
    # PORCENTAJES
    # ----------------------------
    totales = df_agrupado.groupby(col_x)[col_valores].transform("sum")
    df_agrupado["Porcentaje"] = (df_agrupado[col_valores] / totales * 100).round(1)
    # ----------------------------
    # ORDEN CORRECTO (CLAVE)
    # ordenado por TOTAL REAL del valor (ej: Actualizados)
    # ----------------------------
    categorias_color = list(df_agrupado[col_color].unique())
    # depues de este se puede aplicar la lógica para invertir el orden de las categorias en barras
    orden_x = (
        df_agrupado[df_agrupado[col_color]==categorias_color[0]]
        .sort_values("Porcentaje", ascending=orden_ascendente)
    )[col_x].tolist()
    # ----------------------------
    # ORDEN DE CATEGORÍAS (STACK)
    # ----------------------------
    categorias = df_agrupado[col_color].unique().tolist()
    if invertir_apilado:
        categorias = categorias[::-1]

    # ----------------------------
    # PALETA DE COLORES
    # ----------------------------
    if lista_colores:
        color_map = {
            cat: lista_colores[i % len(lista_colores)]
            for i, cat in enumerate(categorias)
        }
    else:
        color_map = {
            cat: ["#235B4E", "#D4C19C", "#621132", "#3A7D6B"][i % 4]
            for i, cat in enumerate(categorias)
        }

    # ----------------------------
    # GRÁFICA
    # ----------------------------
    fig = px.bar(
        df_agrupado,
        x=col_x,
        y="Porcentaje",
        color=col_color,
        text=df_agrupado["Porcentaje"].map("{:.1f}%".format),
        color_discrete_map=color_map,
        category_orders={
            col_x: orden_x,
            col_color: categorias
        },
        barmode="stack",
    )

    # ----------------------------
    # ESTILO DE BARRAS
    # ----------------------------
    fig.update_traces(
        textposition="inside",
        insidetextanchor="middle",
        textfont=dict(size=font_size, color="white"),
        hovertemplate=(
            "<b>%{x}</b><br>"
            "<b>Categoría:</b> %{customdata[1]}<br>"
            "<b>Cantidad:</b> %{customdata[0]:,.0f}<br>"
            "<b>Porcentaje:</b> %{y:.1f}%"
            "<extra></extra>"
        ),
        customdata=df_agrupado[[col_valores, col_color]].values
    )

    # ----------------------------
    # DISEÑO INSTITUCIONAL
    # ----------------------------
    fig.update_layout(
        title=dict(
            text=f"<b>{titulo}</b>",
            font=dict(size=font_size+10, color="#621132")
        ),
        xaxis=dict(
            title="",
            tickangle=-35,
            tickfont=dict(size=font_size, family=FONT_FAMILY,color="black"),
        ),
        yaxis=dict(
            title=dict(
                text="Porcentaje (%)",
                font=dict(
                size=font_size,          # Tamaño del título
                family=FONT_FAMILY,
                color="black"
                ),
            ),
            range=[0, 110],
            ticksuffix="%",
            tickfont=dict(size=font_size, family=FONT_FAMILY,color="black"),
        ),
        legend=dict(
            title=dict(
                text=f"<b>{str(col_color).upper()}</b>",
                font=dict(size=font_size+4, color="black", family=FONT_FAMILY)
            ),
            # y=1.02, 
            orientation="h",
            y=-0.4,
            x=0.5,
            # yanchor="bottom",
            xanchor="center",
            font=dict(
                size=font_size+4,   # <- tamaño de No / Si
                family=FONT_FAMILY,
                color="black"
            ),
        ),
        bargap=0.25,
        height=height,
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(0,0,0,0)",
        margin=dict(l=40, r=20, t=80, b=120),
        # uniformtext_minsize=14,
        # uniformtext_mode="hide",
    )

    return fig

def crear_dona(
    df,
    titulo="",
    colores_lista=None,
    height=580,
    font_size_labels=20,
    hole=0.45,
    add_text_center = True,
):
    """
    Dona con etiquetas internas.

    DataFrame esperado:
        Columna 0 -> Categoría
        Columna 1 -> Valor

    Ejemplos:
        df[['Sexo','Personas']]
        df[['Edad','Personas']]
        df[['Escolaridad','Personas']]
    """

    if df is None or df.empty or len(df.columns) < 2:
        return px.pie(title="Sin datos")

    if colores_lista is None:
        colores_lista = PALETA_INSTITUCIONAL

    df = df.iloc[:, :2].copy()

    col_categoria = df.columns[0]
    col_valor = df.columns[1]
    total = df[col_valor].sum()

    # Etiquetas de 3 líneas
    df["Etiqueta"] = df.apply(
        lambda r: (
            f"<b>{r[col_categoria]}</b><br>"
            f"({r[col_valor]/total:.1%})"
        ),
        axis=1,
    )

    fig = px.pie(
        df,
        values=col_valor,
        names="Etiqueta",
        hole=hole,
        color_discrete_sequence=colores_lista[:len(df)],
    )

    # Total central
    if add_text_center:
        fig.add_annotation(
            x=0.5,
            y=0.5,
            showarrow=False,
            text=f"<b>{total:,.0f}</b><br>personas",
            font=dict(
                size=FONT_SIZE_TITLE,
                color=GUINDA,
                family=FONT_FAMILY
            ),
        )

    fig.update_traces(
        textposition="inside",
        textinfo="value+percent",
        textfont=dict(
            size=font_size_labels,
            color="white",
            family=FONT_FAMILY,
        ),
        # insidetextorientation="horizontal",
        marker=dict(
            line=dict(
                color="white",
                width=2
            )
        ),
        customdata=df[[col_categoria]],
        hovertemplate=(
            "<b>%{customdata[0]}:</b><br>"
            "%{value:,.0f} personas<br>"
            "%{percent}"
            "<extra></extra>"
        ),
    )

    fig.update_layout(
        title=dict(
            text=titulo,
            x=0.5,
            xanchor="center",
            font=dict(
                size=FONT_SIZE_TITLE,
                color=GUINDA,
                family=FONT_FAMILY,
            ),
        ),
        hoverlabel=dict(
            font_size=font_size_labels,
            font_family=FONT_FAMILY,
            bgcolor="white",
            font_color=VERDE,
            bordercolor=DORADO,
        ),
        legend=dict(
            font=dict(
                size=20,
                family=FONT_FAMILY
            ),
            orientation="h",
            yanchor="top",
            y=-0.10,
            xanchor="center",
            x=0.5,
        ),
        height=height,
        margin=dict(
            t=70,
            b=60,
            l=20,
            r=20,
        ),
        showlegend=True,
        paper_bgcolor="rgba(0,0,0,0)",
        uniformtext_minsize=12,
    )

    return fig

def grafica_cumsum(
    df,
    periodo="semana",
    n=5,
    titulo=None,
    color="#3A7D6B",
    fill_color="rgba(58,125,107,0.15)",
    titulo_color="#691C32",
    text_color="#691C32",
    text_size=18,
    height=430,
):

    # Agrupar y ordenar
    df_plot = (
        df.groupby(periodo, as_index=False)["Personas"]
          .sum()
          .sort_values(periodo)
    )

    # Acumulado histórico
    df_plot["Acumulado"] = df_plot["Personas"].cumsum()

    # Últimos n periodos
    df_plot = df_plot.tail(n)
    df_plot["periodo_texto"] = "."+pd.to_datetime(df_plot[periodo]).dt.strftime("%m-%d")
    fig = go.Figure()

    fig.add_trace(
        go.Scatter(
            x=df_plot["periodo_texto"],
            y=df_plot["Acumulado"],
            mode="lines+markers+text",

            text=df_plot["Acumulado"].map("{:,.0f}".format),
            textposition="top center",
            cliponaxis=False,

            textfont=dict(
                size=text_size,
                color=text_color,
                family=FONT_FAMILY
            ),

            line=dict(
                color=color,
                width=3
            ),

            marker=dict(
                size=14,
                color=color,
                line=dict(
                    color="white",
                    width=2
                )
            ),

            fill="tozeroy",
            fillcolor=fill_color,

            customdata=df_plot["Personas"],

            hovertemplate=(
                "<b>%{x}</b><br>"
                "Acumulado: %{y:,.0f}<br>"
                "Periodo: %{customdata:,.0f} personas"
                "<extra></extra>"
            ),
        )
    )

    y_max = df_plot["Acumulado"].max()
    y_min = df_plot["Acumulado"].min()

    y_padding = (y_max - y_min) * 0.25  # 25% de aire arriba


    fig.update_layout(

        title=dict(
            text=titulo or f"Acumulado por {periodo.capitalize()}",
            font=dict(
                size=FONT_SIZE_TITLE,
                color=titulo_color,
                family=FONT_FAMILY
            ),
            x=0
        ),

        height=height,

        template="plotly_white",

        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(0,0,0,0)",

        margin=dict(l=60, r=60, t=45, b=20),

        hoverlabel=dict(
            bgcolor="white",
            font_size=16,
            font_family=FONT_FAMILY
        ),

        xaxis=dict(
            title=" ",
            tickangle=-45,
            showgrid=False,
            zeroline=False,
            showline=False,
            tickfont=dict(
                size=text_size-10,
                family=FONT_FAMILY
            )
        ),

        yaxis=dict(
            title="Personas acumuladas",
            gridcolor="rgba(0,0,0,.08)",
            gridwidth=1,
            range=[y_min, y_max + y_padding],
            autorange=False,
            zeroline=False,
            showline=False,
            visible=False,
            tickfont=dict(
                size=18,
                family=FONT_FAMILY
            )
        ),
    )

    return fig

def add_styled_line(
    text_frame,
    parts,
    font_size=28,
    font_name="Noto Sans"
):
    """
    Crea una sola línea con múltiples estilos (runs).

    Parameters
    ----------
    text_frame : pptx.text.text.TextFrame
        Contenedor del texto (shape.text_frame)

    parts : list of tuples
        [(texto, RGBColor, bold), ...]

    font_size : int
    font_name : str
    """

    text_frame.clear()
    p = text_frame.paragraphs[0]

    for text, color, bold in parts:
        run = p.add_run()
        run.text = text

        run.font.size = Pt(font_size)
        run.font.name = font_name
        run.font.color.rgb = color
        run.font.bold = bold

    return p

def add_biñetas(
    slide,
    df,
    category_col,
    value_col,
    left,
    top,
    width,
    height,
    title="",
    font_name="Noto Sans",
    title_size=24,
    body_size=20,
    title_color=RGB_ROJO,
    value_color=RGB_VERDE,
    text_color=RGB_VERDE_CLARO,
):
    """
    Agrega un cuadro de texto con resumen por categorías.

    Parameters
    ----------
    df : DataFrame
    category_col : str
    value_col : str
    """

    resumen = (
        df.groupby(category_col, as_index=False)[value_col]
          .sum()
          .sort_values(value_col, ascending=False)
    )

    total = resumen[value_col].sum()

    shape = slide.shapes.add_textbox(
        left, top, width, height
    )

    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    # --------------------------
    # Título
    # --------------------------

    p = tf.paragraphs[0]

    p.text = title
    p.alignment = PP_ALIGN.LEFT

    p.font.name = font_name
    p.font.bold = True
    p.font.size = Pt(title_size)
    p.font.color.rgb = title_color

    # --------------------------
    # Categorías
    # --------------------------

    for _, row in resumen.iterrows():

        pct = 100 * row[value_col] / total if total else 0

        p = tf.add_paragraph()

        # Bullet
        run = p.add_run()
        run.text = "• "

        run.font.name = font_name
        run.font.size = Pt(body_size)
        run.font.color.rgb = text_color

        # Valor
        run = p.add_run()
        run.text = f"{row[value_col]:,.0f}"

        run.font.name = font_name
        run.font.bold = True
        run.font.size = Pt(body_size)
        run.font.color.rgb = value_color

        # %
        run = p.add_run()
        run.text = f" ({pct:.1f}%) "

        run.font.name = font_name
        run.font.bold = True
        run.font.size = Pt(body_size)
        run.font.color.rgb = RGBColor(*value_color)

        # Categoría
        run = p.add_run()
        run.text = str(row[category_col])

        run.font.name = font_name
        run.font.size = Pt(body_size)
        run.font.color.rgb = RGBColor(*text_color)

    return shape

def add_kpi(slide, x, y, w, h, titulo, valor, porcentaje):
    VERDE = RGBColor(0x28, 0x5C, 0x4D)
    GUINDA = RGBColor(0x62, 0x11, 0x32)
    VERDE_CLARO = RGBColor(0x3A, 0x7D, 0x6B)
    DORADO = RGBColor(0xD4, 0xC1, 0x9C)

    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        x, y, w, h
    )

    # Fondo transparente
    card.fill.background()

    # Borde dorado
    card.line.color.rgb = DORADO
    card.line.width = Pt(2)

    tf = card.text_frame
    tf.clear()

    # Centrado vertical
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    # Línea 1
    p1 = tf.paragraphs[0]
    p1.text = str(titulo)
    p1.alignment = PP_ALIGN.CENTER

    run1 = p1.runs[0]
    run1.font.size = Pt(18)
    run1.font.bold = True
    run1.font.color.rgb = VERDE

    # Línea 2
    p2 = tf.add_paragraph()
    p2.text = str(valor)
    p2.alignment = PP_ALIGN.CENTER

    run2 = p2.runs[0]
    run2.font.size = Pt(34)
    run2.font.bold = True
    run2.font.color.rgb = GUINDA

    # Línea 3
    p3 = tf.add_paragraph()
    p3.text = str(porcentaje)
    p3.alignment = PP_ALIGN.CENTER

    run3 = p3.runs[0]
    run3.font.size = Pt(16)
    run3.font.bold = True
    run3.font.color.rgb = VERDE_CLARO

    return card

def add_tarjeta_bullets(
    slide,
    df,
    category_col,
    value_col,
    left,
    top,
    width,
    height,
    title="",
    font_name="Noto Sans",
    title_size=18,
    body_size=14,
    header_fill=RGB_VERDE,
    header_color=RGB_BLANCO,
    border_color=RGB_DORADO,
    category_color=RGB_DORADO,
    value_color=RGB_ROJO,
    show_text = True,
):
    """
    Tarjeta institucional con encabezado y cuerpo.

    ┌──────────────────────────────┐
    │         TÍTULO               │
    ├──────────────────────────────┤
    │ • 235 (42.5%) Mujeres        │
    │ • 198 (35.8%) Hombres        │
    │ • 120 (21.7%) Otro           │
    └──────────────────────────────┘
    """

    #===========================
    # Resumen
    #===========================

    resumen = (
        df.groupby(category_col, as_index=False)[value_col]
          .sum()
          .sort_values(value_col, ascending=False)
    )

    total = resumen[value_col].sum()

    #===========================
    # Dimensiones
    #===========================

    header_h = Inches(title_size / 72 * 1.5)
    body_h = height - header_h

    #===========================
    # Encabezado
    #===========================

    header = slide.shapes.add_shape(
        MSO_SHAPE.ROUND_2_SAME_RECTANGLE,
        left,
        top,
        width,
        header_h,
    )

    header.fill.solid()
    header.fill.fore_color.rgb = header_fill
    header.line.fill.rgb = border_color
    header.adjustments[0] = 0.9 

    tf = header.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    run = p.add_run()
    run.text = title
    run.font.name = font_name
    run.font.bold = True
    run.font.size = Pt(title_size)
    run.font.color.rgb = header_color

    #===========================
    # Cuerpo
    #===========================

    body = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left,
        top + header_h,
        width,
        body_h,
    )

    body.fill.solid()
    body.fill.fore_color.rgb = RGB_BLANCO

    body.line.color.rgb = border_color
    body.line.width = Pt(1)

    tf = body.text_frame
    tf.clear()

    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    if show_text:
        for i, (_, row) in enumerate(resumen.iterrows()):

            pct = 100 * row[value_col] / total if total else 0

            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

            p.space_after = Pt(3)
            p.alignment = PP_ALIGN.LEFT

            # Bullet
            r = p.add_run()
            r.text = "• "
            r.font.name = font_name
            r.font.size = Pt(body_size)
            r.font.color.rgb = category_color

            # Valor
            r = p.add_run()
            r.text = f"{row[value_col]:,.0f}"
            r.font.bold = True
            r.font.name = font_name
            r.font.size = Pt(body_size)
            r.font.color.rgb = value_color

            # %
            r = p.add_run()
            r.text = f" ({pct:.1f}%) "
            r.font.bold = True
            r.font.name = font_name
            r.font.size = Pt(body_size)
            r.font.color.rgb = value_color

            # Categoría
            r = p.add_run()
            r.text = str(row[category_col])
            r.font.name = font_name
            r.font.size = Pt(body_size)
            r.font.color.rgb = category_color

    return body

def agregar_forma(
    slide,
    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
    left=None,
    top=None,
    width=None,
    height=None,
    fill_color=RGB_VERDE,
    border_color=RGB_DORADO,
    border_width=1,
    curvature=0.5,
    rotation=0
):
    """
    Agrega una forma geométrica configurando posición, dimensiones, colores y rotación.

    Argumentos:
    - slide: Objeto de la diapositiva activa.
    - shape_type: Tipo de forma (ej: MSO_SHAPE.ROUND_1_CORNER_RECTANGLE).
    - left, top, width, height: Coordenadas y dimensiones (Inches() o Pt()).
    - fill_color: Color de relleno en formato RGBColor(R, G, B).
    - border_color: Opcional. Color del borde en formato RGBColor(R, G, B).
    - border_width: Opcional. Grosor del borde en puntos (Pt). Por defecto es 1.
    - curvature: Opcional. Ajuste de curvatura (float de 0.0 a 1.0) para formas redondeadas.
    - rotation: Opcional. Ángulo de rotación en grados (0 a 360). Por defecto es 0.
    """
    
    # 1. Insertar la forma en la diapositiva
    forma = slide.shapes.add_shape(shape_type, left, top, width, height)
    spTree = slide.shapes._spTree
    spTree.remove(forma._element)
    spTree.insert(2, forma._element)

    # 2. Configurar la rotación (Acepta valores de 0 a 360 grados)
    if rotation != 0:
        forma.rotation = rotation
    
    # 3. Configurar el color de relleno sólido
    forma.fill.solid()
    forma.fill.fore_color.rgb = fill_color
    
    # 4. Configurar el borde o contorno
    if border_color is not None:
        forma.line.color.rgb = border_color
        forma.line.width = Pt(border_width)
    else:
        forma.line.fill.background()  # Quita el contorno por completo
        
    # 5. Ajustar curvatura para formas redondeadas
    if curvature is not None and hasattr(forma, 'adjustments') and len(forma.adjustments) > 0:
        try:
            forma.adjustments = curvature
        except:
            pass
            
    return forma




hoy = datetime.now(ZoneInfo("America/Mexico_City"))
# prs = Presentation("plantilla.pptx")

# df = pl.read_parquet("concentrado_actualizados.parquet").to_pandas()
# print(df.columns)


# #########################################################################################################
# # Titulo de la presentación
# slide_titulo = prs.slides.add_slide(prs.slide_layouts[1])
# # Usa el placeholder de título existente
# slide_titulo.shapes.title.text = "Producción para el Bienestar"
# slide_titulo.placeholders[1].text = "Proceso de actualización - 2026"

fecha = (format_date(hoy, format="d 'de' MMMM 'de' yyyy", locale="es"))
# fecha_datos =format_date(df["dia"].max(), format="d 'de' MMMM 'de' yyyy", locale="es")
# shape_procesos = slide_titulo.shapes.add_textbox(left=Inches(0.25),top=Inches(6.5),width=Inches(5),height=Inches(0.5))
# add_styled_line(shape_procesos.text_frame, [("Reporte con información\n",RGB_BLANCO,True),(f"al {fecha_datos}",RGB_BLANCO,True)], font_size=22)
# #########################################################################################################


# #########################################################################################################
# # Separador de seccion
# slide_h1 = prs.slides.add_slide(prs.slide_layouts[16])
# slide_h1.shapes.title.text = "Avance general"
# # #########################################################################################################

# # #########################################################################################################
# # # Separador de seccion
# slide_S1h1 = prs.slides.add_slide(prs.slide_layouts[11])
# slide_S1h1.shapes.title.text = "Avance Nacional"

# meta = df['Personas'].sum() if len(df) > 0 else 0
# actualizados = df[df['ACTUALIZADO'] == 'Si']['Personas'].sum() if len(df) > 0 else 0
# pendientes = meta - actualizados
# pct_avance = actualizados / meta * 100 if meta > 0 else 0
# # pct_pago = (df[df['Pagados_2026'] > 0 & df['Pagados_2026'].notnull() ]['Personas'].sum() if len(df) > 0 else 0) / actualizados * 100 if actualizados > 0 else 0

# x, y, w, h = Inches(1), Inches(1.1), Inches(3.3), Inches(1.3)
# sep = Inches(4)
# add_kpi(
#     slide_S1h1,
#     x, y, w, h,
#     "Meta de Personas",
#     f"{meta:,.0f}",
#     "100%",
# )

# add_kpi(
#     slide_S1h1,
#     x+sep, y, w, h,
#     "Personas actualizadas",
#     f"{actualizados:,.0f}",
#     f"{pct_avance:.1f}%",
# )

# add_kpi(
#     slide_S1h1,
#     x+2*sep, y, w, h,
#     "Personas pendientes",
#     f"{pendientes:,.0f}",
#     f"{100 - pct_avance:.1f}%",
# )

# barra_estado = crear_barras_porcentaje(df,"NOM_REP", "ACTUALIZADO", "Personas","", orden_ascendente=True, invertir_apilado=True,font_size = 22, height=480)
# # Convertir Plotly → PNG en memoria
# img_bytes = barra_estado.to_image(
#     format="png",
#     width=1880,
#     height=800,
#     scale=1
# )

# img_stream = io.BytesIO(img_bytes)

# # Insertar en PowerPoint
# slide_S1h1.shapes.add_picture(
#     img_stream,
#     left=Inches(0.15),
#     top=Inches(1.9),
#     width=Inches(13)
# )
# # #########################################################################################################
# #########################################################################################################
# # Separador de seccion
# slide_S1h2 = prs.slides.add_slide(prs.slide_layouts[11])
# slide_S1h2.shapes.title.text = "Avance Nacional - OREF"

# resultado = (
#     df.assign(avance=lambda x: x["Personas"].where(x["ACTUALIZADO"] == "Si", 0))
#     .groupby(["CVE_REP_PROD","NOM_REP"], as_index=False)
#     .agg(meta=("Personas", "sum"),avance=("avance", "sum"))
#     .query("avance > 0")
#     .assign(porcentaje=lambda x: (100 * x["avance"] / x["meta"]))
#     .sort_values("porcentaje",ascending=False)
#     .assign(porcentaje=lambda x: (100 * x["avance"] / x["meta"]).map("{:.2f}%".format))
# )
# print(resultado)
# n = len(resultado)
# n_t1 = (n + 1) // 2   # primera mitad (redondea hacia arriba)
# n_t2 = n // 2  

# top_head=Inches(1.3)
# row_head_h=Inches(0.65)
# top_body=top_head+row_head_h


# add_table_body(
#     slide_S1h2,
#     resultado.drop("CVE_REP_PROD", axis=1).head(n_t1),
#     left=Inches(0.5),
#     top=top_body,
#     width=Inches(6),
#     height=Inches(5),
#     font_size=12,
#     column_widths=[30,20,20,20],
#     column_alignments=["left","right","right","right"],
#     row_height=0.55,
# )

# add_header_row(
#     slide=slide_S1h2,
#     left=Inches(0.5),
#     top=top_head,
#     width=Inches(6),
#     text=["OREF","Meta\n(Personas)","Avance\n(Personas)","Avance\n(%)"],
#     header_colors=["#10312B","#691C32","#235B4E","#235B4E"],
#     column_widths=[30,20,20,20],
#     column_alignments=["center","center","center","center"],
#     row_height=0.65,
#     font_size=14,
#     font_color=(255,255,255)
# )

# add_table_body(
#     slide_S1h2,
#     resultado.drop("CVE_REP_PROD", axis=1).tail(n_t2),
#     left=Inches(7),
#     top=top_body,
#     width=Inches(6),
#     height=Inches(5),
#     font_size=12,
#     column_widths=[30,20,20,20],
#     column_alignments=["left","right","right","right"],
#     row_height=0.55,
# )

# add_header_row(
#     slide=slide_S1h2,
#     left=Inches(7),
#     top=top_head,
#     width=Inches(6),
#     text=["OREF","Meta\n(Personas)","Avance\n(Personas)","Avance\n(%)"],
#     header_colors=["#10312B","#691C32","#235B4E","#235B4E"],
#     column_widths=[30,20,20,20],
#     column_alignments=["center","center","center","center"],
#     row_height=0.65,
#     font_size=14,
#     font_color=(255,255,255)
# )


# slide_S1h3 = prs.slides.add_slide(prs.slide_layouts[7])
# slide_S1h3.shapes.title.text = "\t"

# shape_titulo = slide_S1h3.shapes.add_textbox(left=Inches(0.25),top=Inches(0.25),width=Inches(10),height=Inches(2))
# add_styled_line(shape_titulo.text_frame, parts=[(f"{actualizados:,.0f} personas actualizadas\n",RGB_VERDE,True),("con los siguientes indicadores:",RGB_VERDE,True)],font_size=34,font_name="Noto Sans SemiBold")

# COLOR_VALOR = RGB_ROJO
# COLOR_TEXTO = RGB_DORADO


# # agregar_forma(
# #     slide_S1h3,
# #     shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
# #     left=Inches(0.1),
# #     top=Inches(0.1),
# #     width=Inches(8.45),
# #     height=Inches(1.3),
# #     fill_color=RGB_BLANCO,
# #     border_color=RGB_DORADO,
# #     border_width=0.5,
# #     curvature=0.2,
# #     rotation=0,
# # )

# # Imagen nacional
# # slide_S1h3.shapes.add_picture(
# #     f"estados/{int(oref):02d}.png",
# #     Inches(7),
# #     Inches(0.1),
# #     width=Inches(1.3),
# # )

# sep_top = Inches(1.4)
# sep_x = Inches(4.4)
# x, y, w, h = Inches(0.1), Inches(0.1)+sep_top, Inches(4.1), Inches(1.3)


# add_tarjeta_bullets(
#     slide_S1h3,
#     df=df[df['ACTUALIZADO'] == 'Si'],
#     category_col="genero",
#     value_col="Personas",
#     left=x,
#     top=y,
#     width=w,
#     height=h,
#     title="Género",
#     body_size=18,
#     title_size=22,
#     header_fill=RGB_VERDE,
#     category_color=COLOR_TEXTO,
#     value_color=COLOR_VALOR,
# )

# add_tarjeta_bullets(
#     slide_S1h3,
#     df=df[df['ACTUALIZADO'] == 'Si'],
#     category_col="escala",
#     value_col="Personas",
#     left=x+sep_x,
#     top=y,
#     width=w,
#     height=h,
#     title="Escala",
#     body_size=18,
#     title_size=22,
#     header_fill=RGB_ROJO,
#     category_color=COLOR_TEXTO,
#     value_color=COLOR_VALOR,
# )

# add_tarjeta_bullets(
#     slide_S1h3,
#     df=df[df['ACTUALIZADO'] == 'Si'],
#     category_col="Estatus_coordenadas",
#     value_col="Personas",
#     left=x+sep_x*2,
#     top=y,
#     width=w,
#     height=h,
#     title="Estatus coordenadas",
#     body_size=18,
#     title_size=22,
#     header_fill=RGB_CAFE,
#     category_color=COLOR_TEXTO,
#     value_color=COLOR_VALOR,
# )

# x, y, w, h = Inches(0.1), Inches(0.2)+sep_top+h, Inches(4.1), Inches(1.6)

# add_tarjeta_bullets(
#     slide_S1h3,
#     df=df[df['ACTUALIZADO'] == 'Si'],
#     category_col="ciclo",
#     value_col="Personas",
#     left=x,
#     top=y,
#     width=w,
#     height=h,
#     title="Ciclo",
#     body_size=18,
#     title_size=22,
#     header_fill=RGB_VERDE_CLARO,
#     category_color=COLOR_TEXTO,
#     value_color=COLOR_VALOR,
# )

# add_tarjeta_bullets(
#     slide_S1h3,
#     df=df[df['ACTUALIZADO'] == 'Si'],
#     category_col="regimen_predominante",
#     value_col="Personas",
#     left=x+sep_x,
#     top=y,
#     width=w,
#     height=h,
#     title="Régimen hídrico",
#     body_size=18,
#     title_size=22,
#     header_fill=RGB_ROJO_CLARO,
#     category_color=COLOR_TEXTO,
#     value_color=COLOR_VALOR,
# )

# add_tarjeta_bullets(
#     slide_S1h3,
#     df=df[df['ACTUALIZADO'] == 'Si'],
#     category_col="Cambio_predios",
#     value_col="Personas",
#     left=x+sep_x*2,
#     top=y,
#     width=w,
#     height=h,
#     title="Cambios en predios",
#     body_size=18,
#     title_size=22,
#     header_fill=RGB_DORADO,
#     category_color=COLOR_TEXTO,
#     value_color=COLOR_VALOR,
# )

# add_tarjeta_bullets(
#     slide_S1h3,
#     df=df[df['ACTUALIZADO'] == 'Si'],
#     category_col="Estrategia_predominante",
#     value_col="Personas",
#     left=x,
#     top=y+h+Inches(0.1),
#     width=w,
#     height=Inches(2.8),
#     title="Estrategia",
#     body_size=17,
#     title_size=22,
#     header_fill=RGB_GRIS,
#     category_color=COLOR_TEXTO,
#     value_color=COLOR_VALOR,
# )


# add_tarjeta_bullets(
#     slide_S1h3,
#     df=df[df['ACTUALIZADO'] == 'Si'],
#     category_col="escala",
#     value_col="Personas",
#     left=Inches(4.5),
#     top=y+h+Inches(0.1),
#     width=Inches(8.5),
#     height=Inches(2.8),
#     title="Avance acumulado por día",
#     body_size=18,
#     title_size=22,
#     show_text=False,
#     header_fill=RGB_GRIS,
#     category_color=COLOR_TEXTO,
#     value_color=COLOR_VALOR,
# )

# plazo = "dia"
# avance_periodo = grafica_cumsum(df,periodo=plazo,text_size=32,titulo=" ",n=6)

# img_bytes = avance_periodo.to_image(
#     format="png",
#     width=1100,
#     height=480,
#     scale=1
# )

# img_stream = io.BytesIO(img_bytes)

# # Insertar en PowerPoint
# slide_S1h3.shapes.add_picture(
#     img_stream,
#     left=Inches(4.6),
#     top=y+h+Inches(0.5),
#     width=Inches(8.2),
#     height=Inches(2.5),
# )



# #########################################################################################################


# #########################################################################################################
# # Separador para avance por estado
# #########################################################################################################

# #########################################################################################################
# # Separador de seccion
# slide_s2 = prs.slides.add_slide(prs.slide_layouts[16])
# slide_s2.shapes.title.text = "Avance por OREF"
# #########################################################################################################

# #########################################################################################################
# # Contenido por OREF
# orefs = sorted(resultado["CVE_REP_PROD"].dropna().unique())

# for oref in orefs:
    
#     # Filtrar datos de la OREF actual
#     df_oref = df[df["CVE_REP_PROD"] == oref]

#     oref_nombre = df_oref["NOM_REP"].iloc[0] if not df_oref.empty else "N/A"
#     avance_oref = df_oref[df_oref["ACTUALIZADO"] == "Si"]["Personas"].sum()
#     total_personas = df_oref["Personas"].sum()
#     avance_pct_oref = avance_oref / total_personas * 100 if total_personas > 0 else 0

#     # Crear diapositiva
#     slide_s2h1 = prs.slides.add_slide(prs.slide_layouts[17])

#     slide_s2h1.shapes.title.text = (
#         f"{oref_nombre}\n"
#         f"Avance: {avance_oref:,.0f} personas ({avance_pct_oref:.1f}%)"
#     )

#     slide_s2h1.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
#     slide_s2h1.shapes.title.text_frame.paragraphs[1].font.color.rgb = RGBColor(255, 255, 255)
#     slide_s2h1.shapes.title.text_frame.paragraphs[1].font.size = Pt(20)


#     if oref % 3 == 1:
#         COLOR_LAMINA = RGB_VERDE
#         COLOR_VALOR = RGB_ROJO
#         COLOR_TEXTO = RGB_DORADO
#     elif oref % 3 == 2:
#         COLOR_LAMINA = RGB_ROJO
#         COLOR_VALOR = RGB_VERDE
#         COLOR_TEXTO = RGB_DORADO
#     else:
#         COLOR_LAMINA = RGB_DORADO
#         COLOR_VALOR = RGB_ROJO
#         COLOR_TEXTO = RGB_VERDE_CLARO

#     shape_dias_corte = slide_s2h1.shapes.add_textbox(left=Inches(8.8),top=Inches(0.7),width=Inches(4),height=Inches(0.8))
#     add_styled_line(shape_dias_corte.text_frame, [(f"Información al {fecha_datos}\n",COLOR_TEXTO,True),
#                                                 (f"{df_oref.groupby('dia')['Personas'].sum().shape[0]} dias de operación",COLOR_TEXTO,True)], font_size=16)

#     agregar_forma(
#         slide_s2h1,
#         shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
#         left=Inches(0.1),
#         top=Inches(0.1),
#         width=Inches(8.45),
#         height=Inches(1.3),
#         fill_color=COLOR_LAMINA,
#         border_color=RGB_DORADO,
#         border_width=0.5,
#         curvature=0.2,
#         rotation=0,
#     )

#     # Imagen de la OREF
#     slide_s2h1.shapes.add_picture(
#         f"estados/{int(oref):02d}.png",
#         Inches(7),
#         Inches(0.1),
#         width=Inches(1.3),
#     )

#     sep_top = Inches(1.4)
#     sep_x = Inches(4.4)
#     x, y, w, h = Inches(0.1), Inches(0.1)+sep_top, Inches(4.1), Inches(1.3)


#     add_tarjeta_bullets(
#         slide_s2h1,
#         df=df_oref[df_oref['ACTUALIZADO'] == 'Si'],
#         category_col="genero",
#         value_col="Personas",
#         left=x,
#         top=y,
#         width=w,
#         height=h,
#         title="Género",
#         body_size=18,
#         title_size=22,
#         header_fill=COLOR_LAMINA,
#         category_color=COLOR_TEXTO,
#         value_color=COLOR_VALOR,
#     )

#     add_tarjeta_bullets(
#         slide_s2h1,
#         df=df_oref[df_oref['ACTUALIZADO'] == 'Si'],
#         category_col="escala",
#         value_col="Personas",
#         left=x+sep_x,
#         top=y,
#         width=w,
#         height=h,
#         title="Escala",
#         body_size=18,
#         title_size=22,
#         header_fill=COLOR_LAMINA,
#         category_color=COLOR_TEXTO,
#         value_color=COLOR_VALOR,
#     )

#     add_tarjeta_bullets(
#         slide_s2h1,
#         df=df_oref[df_oref['ACTUALIZADO'] == 'Si'],
#         category_col="Estatus_coordenadas",
#         value_col="Personas",
#         left=x+sep_x*2,
#         top=y,
#         width=w,
#         height=h,
#         title="Estatus coordenadas",
#         body_size=18,
#         title_size=22,
#         header_fill=COLOR_LAMINA,
#         category_color=COLOR_TEXTO,
#         value_color=COLOR_VALOR,
#     )

#     x, y, w, h = Inches(0.1), Inches(0.2)+sep_top+h, Inches(4.1), Inches(1.6)

#     add_tarjeta_bullets(
#         slide_s2h1,
#         df=df_oref[df_oref['ACTUALIZADO'] == 'Si'],
#         category_col="ciclo",
#         value_col="Personas",
#         left=x,
#         top=y,
#         width=w,
#         height=h,
#         title="Ciclo",
#         body_size=18,
#         title_size=22,
#         header_fill=COLOR_LAMINA,
#         category_color=COLOR_TEXTO,
#         value_color=COLOR_VALOR,
# )

#     add_tarjeta_bullets(
#         slide_s2h1,
#         df=df_oref[df_oref['ACTUALIZADO'] == 'Si'],
#         category_col="regimen_predominante",
#         value_col="Personas",
#         left=x+sep_x,
#         top=y,
#         width=w,
#         height=h,
#         title="Régimen hídrico",
#         body_size=18,
#         title_size=22,
#         header_fill=COLOR_LAMINA,
#         category_color=COLOR_TEXTO,
#         value_color=COLOR_VALOR,
#     )

#     add_tarjeta_bullets(
#         slide_s2h1,
#         df=df_oref[df_oref['ACTUALIZADO'] == 'Si'],
#         category_col="Cambio_predios",
#         value_col="Personas",
#         left=x+sep_x*2,
#         top=y,
#         width=w,
#         height=h,
#         title="Cambios en predios",
#         body_size=18,
#         title_size=22,
#         header_fill=COLOR_LAMINA,
#         category_color=COLOR_TEXTO,
#         value_color=COLOR_VALOR,
#     )

#     add_tarjeta_bullets(
#         slide_s2h1,
#         df=df_oref[df_oref['ACTUALIZADO'] == 'Si'],
#         category_col="Estrategia_predominante",
#         value_col="Personas",
#         left=x,
#         top=y+h+Inches(0.1),
#         width=w,
#         height=Inches(2.8),
#         title="Estrategia",
#         body_size=17,
#         title_size=22,
#         header_fill=COLOR_LAMINA,
#         category_color=COLOR_TEXTO,
#         value_color=COLOR_VALOR,
#     )


#     add_tarjeta_bullets(
#         slide_s2h1,
#         df=df_oref[df_oref['ACTUALIZADO'] == 'Si'],
#         category_col="escala",
#         value_col="Personas",
#         left=Inches(4.5),
#         top=y+h+Inches(0.1),
#         width=Inches(8.5),
#         height=Inches(2.8),
#         title="Avance acumulado por día",
#         body_size=18,
#         title_size=22,
#         show_text=False,
#         header_fill=COLOR_LAMINA,
#         category_color=COLOR_TEXTO,
#         value_color=COLOR_VALOR,
#     )

#     plazo = "dia"
#     avance_periodo = grafica_cumsum(df_oref,periodo=plazo,text_size=32,titulo=" ",n=6)

#     img_bytes = avance_periodo.to_image(
#         format="png",
#         width=1100,
#         height=480,
#         scale=1
#     )

#     img_stream = io.BytesIO(img_bytes)

#     # Insertar en PowerPoint
#     slide_s2h1.shapes.add_picture(
#         img_stream,
#         left=Inches(4.6),
#         top=y+h+Inches(0.5),
#         width=Inches(8.2),
#         height=Inches(2.5),
#     )

# #########################################################################################################
# # Cierre
# slide_fin = prs.slides.add_slide(prs.slide_layouts[15])
# slide_fin.shapes.title.text = "Gracias"

# #########################################################################################################

# # Guardar archivo
# prs.save(f"{fecha} Reporte Actualización.pptx")

# print("Dashboard creado correctamente.")




# def presentación_asignada_oref(oref,where,parametros):
    #########################################################################################################
    # Contenido por OREF

# username = "eva.orozco"
# oref_asignada = st.secrets["auth"]["credentials"]["usernames"][username].get("oref") if username in st.secrets['auth']['credentials']['usernames'] else []


# where = "WHERE " + " AND ".join([oref_asignada])
@st.cache_data
def descargar_presentacion(oref_asignada= []):

    prs = Presentation("plantilla.pptx")

    oref_con_avance = cargar_datos(f"""
        SELECT DISTINCT "CVE_REP_PROD" FROM concentrado 
        WHERE "CVE_REP_PROD" = ANY(%s) AND "ACTUALIZADO" = 'Si'
        ORDER BY "CVE_REP_PROD";
    """,[oref_asignada])["CVE_REP_PROD"].tolist()

    if not oref_con_avance:
        return

    for oref in oref_con_avance:

        datos_oref = cargar_datos(f"""
        SELECT
            "NOM_REP",
            SUM("Personas") AS "Meta",
            COALESCE(SUM("Personas") FILTER (WHERE "ACTUALIZADO" = 'Si'), 0) AS "Avance",
            ROUND((100.0 * COALESCE(SUM("Personas") FILTER (WHERE "ACTUALIZADO" = 'Si'), 0)/ NULLIF(SUM("Personas"), 0))::numeric,2) AS "Pct"
        FROM concentrado
        WHERE "CVE_REP_PROD" = {oref}
        GROUP BY "NOM_REP";
        """).iloc[0]

        # Filtrar datos de la OREF actual
        oref_nombre, total_personas, avance_oref,  avance_pct_oref = datos_oref
        # Crear diapositiva
        slide_s2h1 = prs.slides.add_slide(prs.slide_layouts[17])
        slide_s2h1.shapes.title.text = (
            f"{oref_nombre}\n"
            f"Avance: {avance_oref:,.0f} personas ({avance_pct_oref:.1f}%)"
        )

        slide_s2h1.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        slide_s2h1.shapes.title.text_frame.paragraphs[1].font.color.rgb = RGBColor(255, 255, 255)
        slide_s2h1.shapes.title.text_frame.paragraphs[1].font.size = Pt(20)


        if oref % 3 == 1:
            COLOR_LAMINA = RGB_VERDE
            COLOR_VALOR = RGB_ROJO
            COLOR_TEXTO = RGB_DORADO
        elif oref % 3 == 2:
            COLOR_LAMINA = RGB_ROJO
            COLOR_VALOR = RGB_VERDE
            COLOR_TEXTO = RGB_DORADO
        else:
            COLOR_LAMINA = RGB_DORADO
            COLOR_VALOR = RGB_ROJO
            COLOR_TEXTO = RGB_VERDE_CLARO

        fecha_datos = cargar_datos(f"""SELECT MAX("dia") FROM concentrado WHERE "CVE_REP_PROD" = {oref};""").iloc[0,0]
        dias_operación = cargar_datos(f"""SELECT COUNT(DISTINCT "dia") FROM concentrado WHERE "CVE_REP_PROD" = {oref};""").iloc[0,0]
        shape_dias_corte = slide_s2h1.shapes.add_textbox(left=Inches(8.8),top=Inches(0.7),width=Inches(4),height=Inches(0.8))
        add_styled_line(shape_dias_corte.text_frame, [(f"Información al {fecha_datos}\n",COLOR_TEXTO,True),
                                                    (f"{dias_operación:0d} dias de operación",COLOR_TEXTO,True)], font_size=16)
        agregar_forma(
            slide_s2h1,
            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
            left=Inches(0.1),
            top=Inches(0.1),
            width=Inches(8.45),
            height=Inches(1.3),
            fill_color=COLOR_LAMINA,
            border_color=RGB_DORADO,
            border_width=0.5,
            curvature=0.2,
            rotation=0,
        )

        # Imagen de la OREF
        slide_s2h1.shapes.add_picture(
            f"estados/{int(oref):02d}.png",
            Inches(7),
            Inches(0.1),
            width=Inches(1.3),
        )

        sep_top = Inches(1.4)
        sep_x = Inches(4.4)
        x, y, w, h = Inches(0.1), Inches(0.1)+sep_top, Inches(4.1), Inches(1.3)

        categoria = "genero"
        df_categoria = cargar_datos(f"""SELECT "{categoria}", sum("Personas") AS "Personas" FROM concentrado WHERE "CVE_REP_PROD" = {oref} AND "ACTUALIZADO"='Si' GROUP BY "{categoria}";""")
        add_tarjeta_bullets(
            slide_s2h1,
            df=df_categoria,
            category_col=categoria,
            value_col="Personas",
            left=x,
            top=y,
            width=w,
            height=h,
            title="Género",
            body_size=18,
            title_size=22,
            header_fill=COLOR_LAMINA,
            category_color=COLOR_TEXTO,
            value_color=COLOR_VALOR,
        )

        categoria = "escala"
        df_categoria = cargar_datos(f"""SELECT "{categoria}", sum("Personas") AS "Personas" FROM concentrado WHERE "CVE_REP_PROD" = {oref} AND "ACTUALIZADO"='Si' GROUP BY "{categoria}";""")
        add_tarjeta_bullets(
            slide_s2h1,
            df=df_categoria,
            category_col=categoria,
            value_col="Personas",
            left=x+sep_x,
            top=y,
            width=w,
            height=h,
            title="Escala",
            body_size=18,
            title_size=22,
            header_fill=COLOR_LAMINA,
            category_color=COLOR_TEXTO,
            value_color=COLOR_VALOR,
        )

        categoria = "Estatus_coordenadas"
        df_categoria = cargar_datos(f"""SELECT "{categoria}", sum("Personas") AS "Personas" FROM concentrado WHERE "CVE_REP_PROD" = {oref} AND "ACTUALIZADO"='Si' GROUP BY "{categoria}";""")
        add_tarjeta_bullets(
            slide_s2h1,
            df=df_categoria,
            category_col=categoria,
            value_col="Personas",
            left=x+sep_x*2,
            top=y,
            width=w,
            height=h,
            title="Estatus coordenadas",
            body_size=18,
            title_size=22,
            header_fill=COLOR_LAMINA,
            category_color=COLOR_TEXTO,
            value_color=COLOR_VALOR,
        )
        x, y, w, h = Inches(0.1), Inches(0.2)+sep_top+h, Inches(4.1), Inches(1.6)

        categoria = "ciclo"
        df_categoria = cargar_datos(f"""SELECT "{categoria}", sum("Personas") AS "Personas" FROM concentrado WHERE "CVE_REP_PROD" = {oref} AND "ACTUALIZADO"='Si' GROUP BY "{categoria}";""")
        add_tarjeta_bullets(
            slide_s2h1,
            df=df_categoria,
            category_col=categoria,
            value_col="Personas",
            left=x,
            top=y,
            width=w,
            height=h,
            title="Ciclo",
            body_size=18,
            title_size=22,
            header_fill=COLOR_LAMINA,
            category_color=COLOR_TEXTO,
            value_color=COLOR_VALOR,
    )

        categoria = "regimen_predominante"
        df_categoria = cargar_datos(f"""SELECT "{categoria}", sum("Personas") AS "Personas" FROM concentrado WHERE "CVE_REP_PROD" = {oref} AND "ACTUALIZADO"='Si' GROUP BY "{categoria}";""")
        add_tarjeta_bullets(
            slide_s2h1,
            df=df_categoria,
            category_col=categoria,
            value_col="Personas",
            left=x+sep_x,
            top=y,
            width=w,
            height=h,
            title="Régimen hídrico",
            body_size=18,
            title_size=22,
            header_fill=COLOR_LAMINA,
            category_color=COLOR_TEXTO,
            value_color=COLOR_VALOR,
        )

        categoria = "Pueblo_originario"
        df_categoria = cargar_datos(f"""SELECT "{categoria}", sum("Personas") AS "Personas" FROM concentrado WHERE "CVE_REP_PROD" = {oref} AND "ACTUALIZADO"='Si' GROUP BY "{categoria}";""")
        add_tarjeta_bullets(
            slide_s2h1,
            df=df_categoria,
            category_col=categoria,
            value_col="Personas",
            left=x+sep_x*2,
            top=y,
            width=w,
            height=h,
            title="Pueblo originario",
            body_size=18,
            title_size=22,
            header_fill=COLOR_LAMINA,
            category_color=COLOR_TEXTO,
            value_color=COLOR_VALOR,
        )

        categoria = "Estrategia_predominante"
        df_categoria = cargar_datos(f"""SELECT "{categoria}", sum("Personas") AS "Personas" FROM concentrado WHERE "CVE_REP_PROD" = {oref} AND "ACTUALIZADO"='Si' GROUP BY "{categoria}";""")
        add_tarjeta_bullets(
            slide_s2h1,
            df=df_categoria,
            category_col=categoria,
            value_col="Personas",
            left=x,
            top=y+h+Inches(0.1),
            width=w,
            height=Inches(2.8),
            title="Estrategia",
            body_size=17,
            title_size=22,
            header_fill=COLOR_LAMINA,
            category_color=COLOR_TEXTO,
            value_color=COLOR_VALOR,
        )


        add_tarjeta_bullets(
            slide_s2h1,
            df=df_categoria,
            category_col=categoria,
            value_col="Personas",
            left=Inches(4.5),
            top=y+h+Inches(0.1),
            width=Inches(8.5),
            height=Inches(2.8),
            title="Avance acumulado por día",
            body_size=18,
            title_size=22,
            show_text=False,
            header_fill=COLOR_LAMINA,
            category_color=COLOR_TEXTO,
            value_color=COLOR_VALOR,
        )

        categoria = "dia"
        df_categoria = cargar_datos(f"""SELECT "{categoria}", sum("Personas") AS "Personas" FROM concentrado WHERE "CVE_REP_PROD" = {oref} AND "ACTUALIZADO"='Si' GROUP BY "{categoria}";""")
        avance_periodo = grafica_cumsum(df_categoria,periodo=categoria,text_size=32,titulo=" ",n=6)
        img_bytes = avance_periodo.to_image(format="png",width=1100,height=480,scale=1)
        img_stream = io.BytesIO(img_bytes)
        # Insertar en PowerPoint
        slide_s2h1.shapes.add_picture(
            img_stream,
            left=Inches(4.6),
            top=y+h+Inches(0.5),
            width=Inches(8.2),
            height=Inches(2.5),
        )

        slide_s2h2 = prs.slides.add_slide(prs.slide_layouts[17])

        slide_s2h2.shapes.title.text = (
            f"{oref_nombre}\n"
            f"Avance: {avance_oref:,.0f} personas ({avance_pct_oref:.1f}%)"
        )

        slide_s2h2.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        slide_s2h2.shapes.title.text_frame.paragraphs[1].font.color.rgb = RGBColor(255, 255, 255)
        slide_s2h2.shapes.title.text_frame.paragraphs[1].font.size = Pt(20)

        fecha_datos = cargar_datos(f"""SELECT MAX("dia") FROM concentrado WHERE "CVE_REP_PROD" = {oref};""").iloc[0,0]
        dias_operación = cargar_datos(f"""SELECT COUNT(DISTINCT "dia") FROM concentrado WHERE "CVE_REP_PROD" = {oref};""").iloc[0,0]
        shape_dias_corte = slide_s2h2.shapes.add_textbox(left=Inches(8.8),top=Inches(0.7),width=Inches(4),height=Inches(0.8))
        add_styled_line(shape_dias_corte.text_frame, [(f"Información al {fecha_datos}\n",COLOR_TEXTO,True),
                                                    (f"{dias_operación:0d} dias de operación",COLOR_TEXTO,True)], font_size=16)
        agregar_forma(
            slide_s2h2,
            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
            left=Inches(0.1),
            top=Inches(0.1),
            width=Inches(8.45),
            height=Inches(1.3),
            fill_color=COLOR_LAMINA,
            border_color=RGB_DORADO,
            border_width=0.5,
            curvature=0.2,
            rotation=0,
        )

        # Imagen de la OREF
        slide_s2h2.shapes.add_picture(
            f"estados/{int(oref):02d}.png",
            Inches(7),
            Inches(0.1),
            width=Inches(1.3),
        )

        sep_top = Inches(1.4)
        sep_x = Inches(4.4)
        x, y, w, h = Inches(0.1), Inches(0.1)+sep_top, Inches(4.1), Inches(1.3)

        df_categoria = pd.DataFrame({'Categoria': ['Personas Afectadas'], 'Personas': [0]})
        add_tarjeta_bullets(
            slide_s2h2,
            df=df_categoria,
            category_col='Categoria',
            value_col="Personas",
            left=x,
            top=y,
            width=Inches(12.9),
            height=Inches(2.7),
            title="Incidencias",
            body_size=18,
            title_size=22,
            header_fill=COLOR_LAMINA,
            category_color=COLOR_TEXTO,
            value_color=COLOR_VALOR,
            show_text=False,
        )

        add_tarjeta_bullets(
            slide_s2h2,
            df=df_categoria,
            category_col='Categoria',
            value_col="Personas",
            left=x,
            top=y+Inches(3.1),
            width=Inches(12.9),
            height=Inches(2.7),
            title="Focos Rojos",
            body_size=18,
            title_size=22,
            header_fill=COLOR_LAMINA,
            category_color=COLOR_TEXTO,
            value_color=COLOR_VALOR,
            show_text=False,
        )
    archivo_pptx = io.BytesIO()
    prs.save(archivo_pptx)
    archivo_pptx.seek(0)
    return archivo_pptx.getvalue()

# descargar_presentacion([1,6,7])
